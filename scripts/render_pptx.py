#!/usr/bin/env python3
"""
PPTX renderer for teacher planners.

Generates child-facing slides with teacher notes for any study node type.
Adapts presentation complexity by key stage:
  - KS1: Large text, simple vocab, lots of emoji, 8-10 slides max
  - KS2: Medium text, enquiry-driven, emoji + text balance, 10-12 slides
  - KS3-4: Denser text, analytical, minimal emoji, 10-14 slides

All content from graph — zero LLM.
"""

import datetime
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from planner_queries import StudyContext

# ── Subject colour palettes ──────────────────────────────────────────

PALETTES = {
    'History':              {'primary': (0xC0, 0x39, 0x2B), 'accent': (0xF3, 0x9C, 0x12), 'bg': (0xFD, 0xF2, 0xE9)},
    'Geography':            {'primary': (0x27, 0xAE, 0x60), 'accent': (0x2E, 0xCC, 0x71), 'bg': (0xEA, 0xFA, 0xF1)},
    'Science':              {'primary': (0x2E, 0x86, 0xC1), 'accent': (0x34, 0x98, 0xDB), 'bg': (0xEB, 0xF5, 0xFB)},
    'English':              {'primary': (0x8E, 0x44, 0xAD), 'accent': (0x9B, 0x59, 0xB6), 'bg': (0xF4, 0xEC, 0xF7)},
    'Mathematics':          {'primary': (0xE6, 0x7E, 0x22), 'accent': (0xF3, 0x9C, 0x12), 'bg': (0xFE, 0xF5, 0xE7)},
    'Art and Design':       {'primary': (0xE7, 0x4C, 0x3C), 'accent': (0xF1, 0x94, 0x8A), 'bg': (0xFD, 0xED, 0xEC)},
    'Music':                {'primary': (0x1A, 0xBC, 0x9C), 'accent': (0x48, 0xC9, 0xB0), 'bg': (0xE8, 0xF8, 0xF5)},
    'Design and Technology': {'primary': (0xD3, 0x54, 0x00), 'accent': (0xE6, 0x7E, 0x22), 'bg': (0xFD, 0xF2, 0xE9)},
    'Computing':            {'primary': (0x2C, 0x3E, 0x50), 'accent': (0x34, 0x49, 0x5E), 'bg': (0xEB, 0xED, 0xEF)},
    'Religious Studies':    {'primary': (0x7D, 0x3C, 0x98), 'accent': (0xAF, 0x7A, 0xC5), 'bg': (0xF4, 0xEC, 0xF7)},
    'Citizenship':          {'primary': (0x16, 0xA0, 0x85), 'accent': (0x1A, 0xBC, 0x9C), 'bg': (0xE8, 0xF6, 0xF3)},
}

# Common colours
DARK_BROWN = RGBColor(0x5D, 0x40, 0x37)
DARK_SLATE = RGBColor(0x2C, 0x3E, 0x50)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
STONE_GREY = RGBColor(0xEC, 0xF0, 0xF1)
CREAM      = RGBColor(0xFD, 0xF2, 0xE9)

# ── Age-band parameters ──────────────────────────────────────────────

AGE_PARAMS = {
    'KS1': {
        'title_size': 48, 'body_size': 28, 'bullet_size': 26, 'note_size': 18,
        'max_slides': 10, 'max_bullet_items': 5, 'emoji_bullets': True,
        'question_prefix': '',
        'vocab_cols': 3, 'vocab_rows': 3,
        'card_font_body': 20, 'card_font_title': 24,
    },
    'KS2': {
        'title_size': 40, 'body_size': 26, 'bullet_size': 22, 'note_size': 18,
        'max_slides': 12, 'max_bullet_items': 6, 'emoji_bullets': True,
        'question_prefix': '',
        'vocab_cols': 4, 'vocab_rows': 3,
        'card_font_body': 20, 'card_font_title': 24,
    },
    'KS3': {
        'title_size': 36, 'body_size': 22, 'bullet_size': 20, 'note_size': 16,
        'max_slides': 14, 'max_bullet_items': 8, 'emoji_bullets': False,
        'question_prefix': 'Key question: ',
        'vocab_cols': 4, 'vocab_rows': 4,
        'card_font_body': 18, 'card_font_title': 22,
    },
    'KS4': {
        'title_size': 36, 'body_size': 20, 'bullet_size': 18, 'note_size': 16,
        'max_slides': 14, 'max_bullet_items': 10, 'emoji_bullets': False,
        'question_prefix': 'Key question: ',
        'vocab_cols': 4, 'vocab_rows': 5,
        'card_font_body': 16, 'card_font_title': 20,
    },
}

SUBJECT_EMOJIS = {
    'History': ['🏛️', '⚔️', '🪙', '📜', '🧱', '🗺️', '👑', '🛡️', '🏺', '🔥'],
    'Geography': ['🌍', '🗺️', '🏔️', '🌊', '🌿', '🏘️', '📊', '🧭', '☀️', '🌧️'],
    'Science': ['🔬', '⚗️', '🧪', '🔭', '💡', '🌡️', '⚡', '🧲', '🦠', '🌱'],
    'English': ['📖', '✍️', '📝', '🎭', '💬', '📚', '🔤', '📃', '🗣️', '💭'],
    'Art and Design': ['🎨', '🖌️', '🖼️', '✂️', '🎭', '📐', '🪡', '🏺', '📸', '🖍️'],
    'Music': ['🎵', '🎶', '🎹', '🎸', '🥁', '🎺', '🎻', '🎤', '🎼', '🔊'],
    'Design and Technology': ['🔧', '🔨', '⚙️', '🪚', '🧵', '📐', '💡', '🔩', '🪛', '🛠️'],
    'Computing': ['💻', '🖥️', '⌨️', '🤖', '📱', '🔗', '🌐', '💾', '🔒', '📡'],
    'Religious Studies': ['🕊️', '📿', '🕌', '⛪', '🕍', '☸️', '🙏', '📖', '🕯️', '🔔'],
    'Citizenship': ['🏛️', '⚖️', '🗳️', '🤝', '📰', '🌍', '💬', '🏠', '📋', '🔑'],
}

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _get_palette(subject: str) -> dict:
    """Get RGB palette for subject, returns dict of RGBColor objects."""
    raw = PALETTES.get(subject, PALETTES.get('History'))
    return {
        'primary': RGBColor(*raw['primary']),
        'accent': RGBColor(*raw['accent']),
        'bg': RGBColor(*raw['bg']),
    }


def _get_age(ks: str) -> dict:
    return AGE_PARAMS.get(ks, AGE_PARAMS['KS2'])


def _emojis(subject: str) -> list[str]:
    return SUBJECT_EMOJIS.get(subject, SUBJECT_EMOJIS['History'])


# ── PPTX helpers ─────────────────────────────────────────────────────

def set_slide_bg(slide, colour):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_text_box(slide, left, top, width, height, text, font_size=24,
                 colour=DARK_BROWN, bold=False, alignment=PP_ALIGN.LEFT,
                 auto_shrink=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if auto_shrink:
        # Enable shrink-text-on-overflow via <a:normAutofit> element
        from pptx.oxml.ns import nsmap
        bodyPr = tf._txBody.bodyPr
        # Remove any existing autofit elements
        for child in list(bodyPr):
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag in ('spAutoFit', 'noAutofit', 'normAutofit'):
                bodyPr.remove(child)
        # Remove any stale attributes from previous broken implementation
        for attr in list(bodyPr.attrib.keys()):
            if 'fontScale' in attr or 'autofit' in attr:
                del bodyPr.attrib[attr]
        # Add normAutofit — allows text to shrink down to 50% to fit the box
        from lxml import etree
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        norm = etree.SubElement(bodyPr, f'{{{ns_a}}}normAutofit')
        norm.set('fontScale', '50000')  # min 50% of original size
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = colour
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=22,
                    colour=DARK_BROWN, emoji_bullets=True, emojis=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    emoji_list = emojis or ['•'] * 20
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if emoji_bullets and emojis:
            bullet = emoji_list[i % len(emoji_list)] + '  '
        else:
            bullet = '•  '
        p.text = bullet + item
        p.font.size = Pt(font_size)
        p.font.color.rgb = colour
        p.font.name = 'Calibri'
        p.space_after = Pt(8)
    return txBox


def add_accent_bar(slide, top, colour, height=Inches(0.08)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), top, Inches(12.333), height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_card(slide, left, top, width, height, bg_colour, border_colour=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_colour
    if border_colour:
        card.line.color.rgb = border_colour
        card.line.width = Pt(2)
    else:
        card.line.fill.background()
    return card


# ── Slide builders ───────────────────────────────────────────────────

def _build_title_slide(prs, ctx, pal, age):
    """Slide 1: Title + hook question."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pal['primary'])

    emoji = _emojis(ctx.subject)[0]
    add_text_box(slide, Inches(5.8), Inches(0.3), Inches(1.5), Inches(1.2),
                 emoji, font_size=80, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.5),
                 ctx.study.get('name', '').upper(),
                 font_size=age['title_size'] + 10, colour=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # First enquiry question as hook
    eqs = ctx.study.get('enquiry_questions', ctx.study.get('enquiry_question', ''))
    if isinstance(eqs, str):
        try:
            eqs = json.loads(eqs)
        except (json.JSONDecodeError, TypeError):
            eqs = [eqs] if eqs else []
    if isinstance(eqs, str):
        eqs = [eqs]

    if eqs:
        add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(1.2),
                     eqs[0], font_size=age['body_size'] + 6, colour=pal['accent'],
                     alignment=PP_ALIGN.CENTER)

    # Study ID line below title
    add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.5),
                 ctx.study_id, font_size=16, colour=RGBColor(0xFF, 0xCC, 0xCC),
                 alignment=PP_ALIGN.CENTER)

    # Metadata bar
    yg = ctx.study.get('year_groups', [])
    if isinstance(yg, list):
        yg = ', '.join(yg)
    period = ctx.study.get('period', '')
    meta = f"{ctx.subject}  •  {ctx.key_stage}"
    if yg:
        meta += f"  •  {yg}"
    if period:
        meta += f"  •  {period}"
    add_text_box(slide, Inches(2), Inches(6.3), Inches(9), Inches(0.6),
                 meta, font_size=18, colour=RGBColor(0xFF, 0xCC, 0xCC),
                 alignment=PP_ALIGN.CENTER)

    # Teacher notes
    notes = [f"TITLE SLIDE — {ctx.study.get('name', '')}"]
    notes.append(f"Study ID: {ctx.study_id}")
    notes.append(f"Generated: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}")
    if eqs:
        notes.append('\nENQUIRY QUESTIONS:')
        for i, q in enumerate(eqs, 1):
            notes.append(f"{i}. {q}")
    pitfalls = ctx.study.get('common_pitfalls', [])
    if isinstance(pitfalls, str):
        try:
            pitfalls = json.loads(pitfalls)
        except (json.JSONDecodeError, TypeError):
            pitfalls = [pitfalls] if pitfalls else []
    if pitfalls:
        notes.append('\nPITFALLS:')
        for p in pitfalls:
            notes.append(f"- {p}")
    sensitive = ctx.study.get('sensitive_content_notes', [])
    if isinstance(sensitive, str):
        try:
            sensitive = json.loads(sensitive)
        except (json.JSONDecodeError, TypeError):
            sensitive = [sensitive] if sensitive else []
    if sensitive:
        notes.append('\nSENSITIVE CONTENT:')
        for s in sensitive:
            notes.append(f"- {s}")
    if ctx.templates:
        vt = ctx.templates[0]
        phases = vt.get('session_structure', [])
        if phases:
            notes.append(f"\nSESSION STRUCTURE ({vt.get('name', '')}):")
            notes.append(' → '.join(phases))
    # Scaffolding guidance from learner profile
    if ctx.learner_profile:
        notes.append('\nSCAFFOLDING GUIDANCE:')
        content = ctx.learner_profile.get('content', {})
        pedagogy = ctx.learner_profile.get('pedagogy', {})
        if content:
            reading_label = content.get('label', '')
            if reading_label:
                notes.append(f"  Reading level: {reading_label}")
            tts = content.get('tts_available', False)
            notes.append(f"  TTS available: {'Yes' if tts else 'No'}")
            lexile_min = content.get('lexile_min')
            lexile_max = content.get('lexile_max')
            if lexile_min is not None and lexile_max is not None:
                notes.append(f"  Lexile range: {lexile_min}–{lexile_max}")
        if pedagogy:
            scaff = pedagogy.get('scaffolding_level', '')
            if scaff:
                notes.append(f"  Scaffolding level: {scaff}")
            sess_min = pedagogy.get('session_length_min_minutes')
            sess_max = pedagogy.get('session_length_max_minutes')
            if sess_min is not None and sess_max is not None:
                notes.append(f"  Session length: {sess_min}–{sess_max} minutes")
            hint_tiers = pedagogy.get('hint_tiers_max')
            if hint_tiers is not None:
                notes.append(f"  Max hint tiers: {hint_tiers}")
    set_notes(slide, '\n'.join(notes))


def _build_concepts_slide(prs, ctx, pal, age):
    """Slide 2: Key concepts as cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pal['bg'])

    emoji = _emojis(ctx.subject)[1]
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 f'{emoji}  Key concepts', font_size=age['title_size'],
                 colour=DARK_BROWN, bold=True)
    add_accent_bar(slide, Inches(1.2), pal['primary'])

    # Show top concepts as cards (max 4-6)
    primary = [c for c in ctx.concepts if c.get('is_primary')]
    secondary = [c for c in ctx.concepts if not c.get('is_primary')]
    show = primary + secondary
    max_cards = 6 if age.get('max_bullet_items', 6) >= 6 else 4
    show = show[:max_cards]

    if not show:
        add_text_box(slide, Inches(1), Inches(2), Inches(10), Inches(1),
                     'No concepts linked to this study yet.',
                     font_size=age['body_size'], colour=DARK_BROWN)
        set_notes(slide, 'No DELIVERS_VIA relationships found for this study node.')
        return

    cols = min(len(show), 3)
    rows = (len(show) + cols - 1) // cols
    card_w = Inches(3.8) if cols == 3 else Inches(5.5) if cols == 2 else Inches(11)
    card_h = Inches(2.3) if rows <= 2 else Inches(1.8)

    emojis = _emojis(ctx.subject)
    for i, c in enumerate(show):
        col = i % cols
        row = i // cols
        left = Inches(0.5 + col * (12.333 / cols))
        top = Inches(1.5 + row * (card_h.inches + 0.3))

        is_pri = c.get('is_primary')
        bg_col = pal['primary'] if is_pri else WHITE
        text_col = WHITE if is_pri else DARK_BROWN
        add_card(slide, left, top, card_w, card_h, bg_col,
                 border_colour=pal['primary'] if not is_pri else None)

        name = c.get('concept_name', c.get('name', ''))
        label = 'Primary' if is_pri else 'Secondary'
        add_text_box(slide, left + Inches(0.15), top + Inches(0.1),
                     card_w - Inches(0.3), Inches(0.5),
                     f"{emojis[i % len(emojis)]}  {name}",
                     font_size=age['card_font_title'],
                     colour=text_col, bold=True)

        desc = c.get('description', '')
        if desc:
            # Truncate for cards
            max_len = 120 if age['body_size'] >= 24 else 180
            if len(desc) > max_len:
                desc = desc[:max_len].rsplit(' ', 1)[0] + '...'
            add_text_box(slide, left + Inches(0.15), top + Inches(0.65),
                         card_w - Inches(0.3), card_h - Inches(0.8),
                         desc, font_size=age['card_font_body'], colour=text_col)

    # Teacher notes with differentiation
    notes = ['CONCEPTS SLIDE\n']
    for c in ctx.concepts:
        cid = c.get('concept_id', '')
        name = c.get('concept_name', c.get('name', ''))
        pri = '(PRIMARY)' if c.get('is_primary') else '(secondary)'
        tw = c.get('teaching_weight', '')
        notes.append(f"{name} {cid} {pri} weight:{tw}/6")
        dls = c.get('difficulty_levels', [])
        if dls:
            notes.append('DIFFERENTIATION:')
            for dl in dls:
                label = dl.get('label', '').replace('_', ' ').title()
                desc = dl.get('description', '')
                notes.append(f"  {label}: {desc}")
                example_resp = dl.get('example_response', '')
                if example_resp:
                    notes.append(f"  Model response ({label}): {example_resp}")
            errors = dls[0].get('common_errors', '')
            if isinstance(errors, list):
                errors = '; '.join(errors)
            if errors:
                notes.append(f"  Common errors (Entry): {errors}")
        notes.append('')
    # KS2 assessment codes
    _add_assessment_codes_to_notes(notes, ctx)
    set_notes(slide, '\n'.join(notes))


def _build_enquiry_slide(prs, ctx, pal, age):
    """Slide 3: Enquiry question + think prompt."""
    eqs = ctx.study.get('enquiry_questions', ctx.study.get('enquiry_question', ''))
    if isinstance(eqs, str):
        try:
            eqs = json.loads(eqs)
        except (json.JSONDecodeError, TypeError):
            eqs = [eqs] if eqs else []
    if isinstance(eqs, str):
        eqs = [eqs]
    if not eqs:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pal['bg'])

    emoji = _emojis(ctx.subject)[2]
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 f'{emoji}  Our big questions', font_size=age['title_size'],
                 colour=DARK_BROWN, bold=True)
    add_accent_bar(slide, Inches(1.2), pal['accent'])

    y = Inches(1.8)
    for i, q in enumerate(eqs[:4]):
        add_text_box(slide, Inches(1), y, Inches(11), Inches(1),
                     f"{i+1}.  {q}", font_size=age['body_size'],
                     colour=DARK_SLATE)
        y += Inches(1.2)

    # Think box
    if ctx.thinking_lenses:
        lens = ctx.thinking_lenses[0]
        kq = lens.get('key_question', '')
        box = add_card(slide, Inches(1), Inches(5.2), Inches(11), Inches(1.8),
                       pal['primary'])
        add_text_box(slide, Inches(1.3), Inches(5.4), Inches(10.5), Inches(1.5),
                     f"Thinking lens: {kq}", font_size=age['body_size'] - 4,
                     colour=WHITE, bold=True)

    notes = ['ENQUIRY QUESTIONS SLIDE']
    if ctx.thinking_lenses:
        lens = ctx.thinking_lenses[0]
        notes.append(f"\nTHINKING LENS: {lens.get('lens_name', '')}")
        notes.append(f"Key question: {lens.get('key_question', '')}")
        if lens.get('rationale'):
            notes.append(f"Rationale: {lens['rationale']}")
        stems = lens.get('question_stems', [])
        if stems:
            notes.append(f"\nQuestion stems ({ctx.key_stage}):")
            for s in stems:
                notes.append(f"- {s}")
    set_notes(slide, '\n'.join(notes))


def _as_list(val):
    """Ensure value is a list."""
    if val is None:
        return []
    if isinstance(val, str):
        try:
            parsed = json.loads(val)
            if isinstance(parsed, list):
                return parsed
        except (json.JSONDecodeError, TypeError):
            pass
        return [val] if val else []
    return val


def _build_subject_content_slide(prs, ctx, pal, age):
    """Slide 4-5: Subject-specific content cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_SLATE)
    study = ctx.study
    emojis = _emojis(ctx.subject)

    # Content depends on subject
    items = []
    title = 'Key content'
    notes_lines = []

    if ctx.subject == 'History':
        title = 'Key people and events'
        figures = _as_list(study.get('key_figures'))
        events = _as_list(study.get('key_events'))
        for f in figures[:4]:
            items.append(('👤', f, ''))
        for e in events[:4]:
            items.append(('📅', e, ''))
        # Rich teacher notes
        notes_lines.append(f"KEY FIGURES: {', '.join(figures)}" if figures else '')
        notes_lines.append(f"KEY EVENTS: {', '.join(events)}" if events else '')
        perspectives = _as_list(study.get('perspectives'))
        if perspectives:
            notes_lines.append(f"\nPERSPECTIVES TO INCLUDE: {', '.join(perspectives)}")
        sig = study.get('significance_claim', '')
        if sig:
            notes_lines.append(f"\nSIGNIFICANCE: {sig}")
        interps = _as_list(study.get('interpretations'))
        if interps:
            notes_lines.append('\nHISTORIOGRAPHICAL DEBATE:')
            for interp in interps:
                notes_lines.append(f"- {interp}")

    elif ctx.subject == 'Science':
        title = 'Investigation setup'
        variables = study.get('variables')
        if isinstance(variables, str):
            try:
                variables = json.loads(variables)
            except (json.JSONDecodeError, TypeError):
                variables = None
        if variables and isinstance(variables, dict):
            items.append(('🔬', f"Independent: {variables.get('independent', '')}", ''))
            items.append(('📏', f"Dependent: {variables.get('dependent', '')}", ''))
            controlled = variables.get('controlled', [])
            if controlled:
                if isinstance(controlled, list):
                    items.append(('🔒', f"Controlled: {', '.join(controlled[:3])}", ''))
        equip = _as_list(study.get('equipment'))
        if equip:
            items.append(('🧪', f"Equipment: {', '.join(equip[:5])}", ''))
        safety = study.get('safety_notes', '')
        if safety:
            items.append(('⚠️', f"Safety: {safety[:60]}", ''))
        outcome = study.get('expected_outcome', '')
        if outcome:
            items.append(('📊', f"Expected: {outcome[:60]}", ''))
        recording = _as_list(study.get('recording_format'))
        if recording:
            items.append(('📝', f"Recording: {', '.join(recording)}", ''))
        # Rich teacher notes
        notes_lines.append(f"ENQUIRY: {study.get('enquiry_question', '')}")
        if safety:
            notes_lines.append(f"\nSAFETY: {safety}")
            hazard = study.get('hazard_level', '')
            if hazard:
                notes_lines.append(f"Hazard level: {hazard}")
        if outcome:
            notes_lines.append(f"\nEXPECTED OUTCOME: {outcome}")
        ped = study.get('pedagogical_rationale', '')
        if ped:
            notes_lines.append(f"\nPEDAGOGICAL RATIONALE: {ped}")
        disc = study.get('science_discipline', '')
        if disc:
            notes_lines.append(f"Science discipline: {disc}")
        ets = ctx.references.get('enquiry_types', [])
        if ets:
            notes_lines.append(f"\nENQUIRY TYPE: {', '.join(et.get('name', '') for et in ets)}")

    elif ctx.subject == 'English':
        title = 'Text study'
        text_type = study.get('text_type', '')
        if text_type:
            items.append(('📖', f"Text type: {text_type.replace('_', ' ').title()}", ''))
        outcome = study.get('writing_outcome', '')
        if outcome:
            items.append(('✍️', f"Writing outcome: {outcome[:70]}", ''))
        grammar = _as_list(study.get('grammar_focus'))
        if grammar:
            items.append(('🔤', f"Grammar: {', '.join(grammar[:4])}", ''))
        features = _as_list(study.get('text_features_to_teach'))
        if features:
            items.append(('📋', f"Features: {', '.join(features[:4])}", ''))
        texts = study.get('suggested_texts', [])
        if isinstance(texts, str):
            try:
                texts = json.loads(texts)
            except (json.JSONDecodeError, TypeError):
                texts = []
        for t in texts[:2]:
            if isinstance(t, dict):
                items.append(('📚', f"{t.get('title', '')} by {t.get('author', '')}", ''))
        lit_terms = _as_list(study.get('literary_terms'))
        if lit_terms:
            items.append(('💡', f"Literary terms: {', '.join(lit_terms[:4])}", ''))
        # Teacher notes
        gram_src = study.get('grammar_year_source', '')
        if gram_src:
            notes_lines.append(f"Grammar from: {gram_src} Appendix 2")
        genres = ctx.references.get('genres', [])
        if genres:
            notes_lines.append(f"\nGENRE: {', '.join(g.get('name', '') for g in genres)}")
            for g in genres:
                if g.get('description'):
                    notes_lines.append(f"  {g['description'][:150]}")

    elif ctx.subject == 'Art and Design':
        title = 'Art focus'
        artist = study.get('artist', '')
        if artist:
            dates = study.get('artist_dates', '')
            items.append(('🎨', f"Artist: {artist}" + (f" ({dates})" if dates else ''), ''))
        movement = study.get('art_movement', '')
        if movement:
            items.append(('🖼️', f"Movement: {movement}", ''))
        medium = _as_list(study.get('medium'))
        if medium:
            items.append(('🖌️', f"Medium: {', '.join(medium)}", ''))
        techniques = _as_list(study.get('techniques'))
        if techniques:
            items.append(('✂️', f"Techniques: {', '.join(techniques[:4])}", ''))
        elements = _as_list(study.get('visual_elements'))
        if elements:
            items.append(('👁️', f"Elements: {', '.join(elements[:4])}", ''))
        context = study.get('cultural_context', '')
        if context:
            items.append(('🌍', f"Context: {context}", ''))

    elif ctx.subject == 'Music':
        title = 'Music focus'
        genre = study.get('genre', '')
        if genre:
            items.append(('🎵', f"Genre: {genre.replace('_', ' ').title()}", ''))
        composer = study.get('composer', '')
        piece = study.get('piece', '')
        if composer:
            items.append(('🎹', f"{'Composer' if composer != piece else 'Piece'}: {composer}", piece if piece != composer else ''))
        elements = _as_list(study.get('musical_elements'))
        if elements:
            items.append(('🔊', f"Elements: {', '.join(elements[:4])}", ''))
        instruments = _as_list(study.get('instrument'))
        if instruments:
            items.append(('🎸', f"Instruments: {', '.join(instruments)}", ''))
        notation = study.get('notation_level', '')
        if notation:
            items.append(('🎼', f"Notation: {notation.replace('_', ' ')}", ''))
        repertoire = _as_list(study.get('listening_repertoire'))
        if repertoire:
            items.append(('👂', f"Listen: {', '.join(repertoire[:3])}", ''))

    elif ctx.subject == 'Design and Technology':
        title = 'Design challenge'
        brief = study.get('design_brief', '')
        if brief:
            items.append(('💡', brief[:90], ''))
        materials = _as_list(study.get('materials'))
        if materials:
            items.append(('🧱', f"Materials: {', '.join(materials[:5])}", ''))
        tools = _as_list(study.get('tools'))
        if tools:
            items.append(('🔧', f"Tools: {', '.join(tools[:5])}", ''))
        techniques = _as_list(study.get('techniques'))
        if techniques:
            items.append(('⚙️', f"Techniques: {', '.join(techniques[:4])}", ''))
        safety = study.get('safety_notes', '')
        if safety:
            items.append(('⚠️', f"Safety: {safety[:70]}", ''))
        eval_criteria = _as_list(study.get('evaluation_criteria'))
        if eval_criteria:
            notes_lines.append('EVALUATION CRITERIA:')
            for e in eval_criteria:
                notes_lines.append(f"- {e}")

    elif ctx.subject == 'Computing':
        title = 'Computing focus'
        tool = study.get('software_tool', '')
        if tool:
            items.append(('💻', f"Tool: {tool}", ''))
        paradigm = study.get('programming_paradigm', '')
        if paradigm:
            items.append(('🤖', f"Paradigm: {paradigm.replace('_', ' ').title()}", ''))
        concepts = _as_list(study.get('computational_concept'))
        if concepts:
            items.append(('⚙️', f"Concepts: {', '.join(c.replace('_', ' ') for c in concepts)}", ''))
        abstraction = study.get('abstraction_level', '')
        if abstraction:
            items.append(('🔍', f"Abstraction: {abstraction.replace('_', ' ').title()}", ''))

    else:
        # Generic fallback — pull whatever properties are interesting
        themes = _as_list(study.get('themes'))
        for t in themes[:4]:
            items.append((emojis[0], t, ''))
        # Also show pedagogical rationale
        ped = study.get('pedagogical_rationale', '')
        if ped:
            items.append(('💡', ped[:80], ''))

    if not items:
        return  # Skip empty slide

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 f'{emojis[3]}  {title}', font_size=age['title_size'],
                 colour=WHITE, bold=True)

    # Render as cards — use auto_shrink for long text
    max_items = min(len(items), 8)
    cols = min(max_items, 4) if max_items > 2 else max_items
    if cols == 0:
        cols = 1
    rows = (max_items + cols - 1) // cols
    card_w = Inches(12 / cols - 0.3)
    card_h = Inches(2.2) if rows <= 2 else Inches(1.6)

    card_colours = [pal['primary'], RGBColor(0x8E, 0x44, 0xAD),
                    RGBColor(0xE6, 0x7E, 0x22), RGBColor(0x2E, 0x86, 0xC1),
                    pal['accent'], RGBColor(0x27, 0xAE, 0x60),
                    RGBColor(0xD3, 0x54, 0x00), RGBColor(0x1A, 0xBC, 0x9C)]

    for i, (emoji, text, sub) in enumerate(items[:max_items]):
        col = i % cols
        row = i // cols
        left = Inches(0.3 + col * (12.5 / cols))
        top = Inches(1.5 + row * (card_h.inches + 0.3))

        bg = card_colours[i % len(card_colours)]
        add_card(slide, left, top, card_w, card_h, bg)

        # Single text frame with title + body as paragraphs
        # This prevents overlap — paragraphs flow naturally within one frame
        txBox = slide.shapes.add_textbox(
            left + Inches(0.15), top + Inches(0.1),
            card_w - Inches(0.3), card_h - Inches(0.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        # Enable shrink-to-fit via normAutofit
        from lxml import etree
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        bodyPr = tf._txBody.bodyPr
        for child in list(bodyPr):
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag in ('spAutoFit', 'noAutofit', 'normAutofit'):
                bodyPr.remove(child)
        norm = etree.SubElement(bodyPr, f'{{{ns_a}}}normAutofit')
        norm.set('fontScale', '50000')

        # Title paragraph (bold, larger)
        p_title = tf.paragraphs[0]
        max_chars = int(card_w.inches * 14)
        display_text = text[:max_chars] + ('...' if len(text) > max_chars else '')
        p_title.text = f"{emoji}  {display_text}"
        p_title.font.size = Pt(age['card_font_title'])
        p_title.font.color.rgb = WHITE
        p_title.font.bold = True
        p_title.font.name = 'Calibri'
        p_title.space_after = Pt(4)

        # Body paragraph (smaller, not bold) — overflow text or subtitle
        remaining = text[max_chars:] if len(text) > max_chars else sub
        if remaining:
            p_body = tf.add_paragraph()
            p_body.text = remaining[:150]
            p_body.font.size = Pt(age['card_font_body'])
            p_body.font.color.rgb = WHITE
            p_body.font.bold = False
            p_body.font.name = 'Calibri'

    # Build comprehensive teacher notes
    notes = [f"SUBJECT CONTENT SLIDE — {title}\n"]
    notes.extend(notes_lines)
    # Add cross-curricular if available
    if ctx.cross_curricular:
        notes.append('\nCROSS-CURRICULAR LINKS:')
        for cc in ctx.cross_curricular:
            subj = cc.get('target_subject', '')
            name = cc.get('target_name', '')
            hook = cc.get('hook', '')
            notes.append(f"- {subj}: {name} — {hook}")
    set_notes(slide, '\n'.join(n for n in notes if n))


def _build_vocabulary_slide(prs, ctx, pal, age):
    """Vocabulary word mat slide."""
    definitions = ctx.study.get('definitions', [])
    if isinstance(definitions, str):
        try:
            definitions = json.loads(definitions)
        except (json.JSONDecodeError, TypeError):
            definitions = []
    if not definitions:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pal['bg'])

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 '📖  Key words', font_size=age['title_size'],
                 colour=DARK_BROWN, bold=True)
    add_accent_bar(slide, Inches(1.0), RGBColor(0x8E, 0x44, 0xAD))

    cols = age['vocab_cols']
    max_items = cols * age['vocab_rows']
    vocab = definitions[:max_items]

    accent_colours = [
        pal['primary'], RGBColor(0x8E, 0x44, 0xAD), RGBColor(0x2E, 0x86, 0xC1),
        RGBColor(0xE6, 0x7E, 0x22), RGBColor(0x27, 0xAE, 0x60), pal['accent'],
    ]

    for i, item in enumerate(vocab):
        col = i % cols
        row = i // cols
        left = Inches(0.3 + col * (12.5 / cols))
        top = Inches(1.3 + row * 2.0)
        card_w = Inches(12 / cols - 0.2)

        # Handle both string and dict definitions
        if isinstance(item, dict):
            word = item.get('term', item.get('word', ''))
            defn = item.get('meaning', item.get('definition', ''))
        else:
            word = str(item)
            defn = ''

        accent = accent_colours[i % len(accent_colours)]
        add_card(slide, left, top, card_w, Inches(1.7), WHITE, border_colour=accent)

        add_text_box(slide, left + Inches(0.1), top + Inches(0.05),
                     card_w - Inches(0.2), Inches(0.45),
                     word, font_size=age['card_font_title'] - 2,
                     colour=accent, bold=True)

        if defn:
            add_text_box(slide, left + Inches(0.1), top + Inches(0.55),
                         card_w - Inches(0.2), Inches(1.05),
                         defn, font_size=age['card_font_body'] - 2,
                         colour=DARK_BROWN)

    notes = ['VOCABULARY SLIDE — Print as word mat for desks.']
    notes.append(f'Terms: {", ".join(str(v) if isinstance(v, str) else v.get("term", v.get("word", "")) for v in vocab)}')
    set_notes(slide, '\n'.join(notes))


def _build_assessment_slide(prs, ctx, pal, age):
    """Assessment / big question slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pal['primary'])

    add_text_box(slide, Inches(5.8), Inches(0.3), Inches(1.5), Inches(1.2),
                 '✍️', font_size=72, alignment=PP_ALIGN.CENTER)

    ks_verbs = {
        'KS1': 'Show what you know!',
        'KS2': 'Your turn to be a historian!' if ctx.subject == 'History' else 'Show what you have learned!',
        'KS3': 'Assessment task',
        'KS4': 'Assessment task',
    }
    heading = ks_verbs.get(ctx.key_stage, 'Assessment task')
    if ctx.subject != 'History':
        heading = heading.replace('historian', 'expert')

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
                 heading, font_size=age['title_size'] + 4,
                 colour=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Assessment question from template or enquiry
    assess_q = ''
    if ctx.templates:
        assess_q = ctx.templates[0].get('assessment_approach', '')
    if not assess_q:
        eqs = ctx.study.get('enquiry_questions', ctx.study.get('enquiry_question', ''))
        if isinstance(eqs, str):
            try:
                eqs = json.loads(eqs)
            except (json.JSONDecodeError, TypeError):
                eqs = [eqs] if eqs else []
        if isinstance(eqs, list) and eqs:
            assess_q = eqs[-1]  # Last enquiry question often best for assessment

    if assess_q:
        add_text_box(slide, Inches(1.5), Inches(3), Inches(10), Inches(1.5),
                     assess_q, font_size=age['body_size'] + 4,
                     colour=pal['accent'], bold=True,
                     alignment=PP_ALIGN.CENTER)

    # Success criteria hint
    success = ctx.study.get('success_criteria', [])
    if isinstance(success, str):
        try:
            success = json.loads(success)
        except (json.JSONDecodeError, TypeError):
            success = []
    if success:
        hints = success[:3]
        y = Inches(4.8)
        for h in hints:
            add_text_box(slide, Inches(2), y, Inches(9), Inches(0.5),
                         f"  •  {h}", font_size=age['bullet_size'],
                         colour=WHITE)
            y += Inches(0.6)

    # Teacher notes with differentiation
    notes = ['ASSESSMENT SLIDE\n']
    if ctx.templates:
        notes.append(f"Assessment approach: {ctx.templates[0].get('assessment_approach', '')}")
    notes.append('\nDIFFERENTIATION BY LEVEL:')
    for c in ctx.concepts[:2]:
        dls = c.get('difficulty_levels', [])
        if dls:
            name = c.get('concept_name', c.get('name', ''))
            notes.append(f"\n{name}:")
            for dl in dls:
                label = dl.get('label', '').replace('_', ' ').title()
                task = dl.get('example_task', dl.get('description', ''))
                notes.append(f"  {label}: {task}")
                example_resp = dl.get('example_response', '')
                if example_resp:
                    notes.append(f"  Model response ({label}): {example_resp}")
    set_notes(slide, '\n'.join(notes))


def _build_next_slide(prs, ctx, pal, age):
    """What's next / transition slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pal['accent'])

    if ctx.leads_to:
        add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
                     "What's next?", font_size=age['title_size'] + 4,
                     colour=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        emoji = _emojis(ctx.subject)[5]
        add_text_box(slide, Inches(5.8), Inches(3), Inches(1.5), Inches(1.2),
                     emoji, font_size=80, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(1.5),
                     f'Next: {ctx.leads_to}',
                     font_size=age['body_size'] + 4, colour=WHITE,
                     alignment=PP_ALIGN.CENTER)
    else:
        add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(2),
                     f"End of {ctx.study.get('name', '')}",
                     font_size=age['title_size'] + 4,
                     colour=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    notes = ['TRANSITION SLIDE']
    if ctx.follows:
        notes.append(f'Follows: {ctx.follows}')
    if ctx.leads_to:
        notes.append(f'Leads to: {ctx.leads_to}')
    if ctx.thinking_lenses:
        notes.append(f"\nThinking lens bridge: {ctx.thinking_lenses[0].get('key_question', '')}")
    set_notes(slide, '\n'.join(notes))


# ── Thinking lens slide ──────────────────────────────────────────────

def _build_thinking_lens_slide(prs, ctx, pal, age):
    """Thinking lens slide — cognitive framing for the study."""
    if not ctx.thinking_lenses:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(0x7C, 0x3A, 0xED))  # ThinkingLens purple

    primary_lens = ctx.thinking_lenses[0]
    lens_name = primary_lens.get('lens_name', '')
    key_question = primary_lens.get('key_question', '')
    question_stems = primary_lens.get('question_stems', [])
    rationale = primary_lens.get('rationale', '')

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 f'🔍  Thinking lens: {lens_name}', font_size=age['title_size'],
                 colour=WHITE, bold=True)

    # Key question in accent box
    add_card(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(1.4),
             RGBColor(0x6C, 0x2B, 0xD9))
    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.2),
                 key_question, font_size=age['body_size'] + 6,
                 colour=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                 auto_shrink=True)

    # Question stems for the key stage
    if question_stems:
        y = Inches(3.0)
        add_text_box(slide, Inches(0.7), y, Inches(11), Inches(0.5),
                     f'Question stems for {ctx.key_stage}:',
                     font_size=age['body_size'] - 2, colour=RGBColor(0xD9, 0xC5, 0xF7))
        y += Inches(0.6)
        for stem in question_stems[:6]:
            add_text_box(slide, Inches(1), y, Inches(11), Inches(0.5),
                         f'•  {stem}', font_size=age['bullet_size'],
                         colour=WHITE, auto_shrink=True)
            y += Inches(0.55)

    # Secondary lens hint
    if len(ctx.thinking_lenses) > 1:
        sec = ctx.thinking_lenses[1]
        add_text_box(slide, Inches(0.7), Inches(6.5), Inches(11), Inches(0.5),
                     f"Alternative lens: {sec.get('lens_name', '')} — {sec.get('key_question', '')}",
                     font_size=age['body_size'] - 4,
                     colour=RGBColor(0xD9, 0xC5, 0xF7))

    # Teacher notes
    notes = [f"THINKING LENS: {lens_name}"]
    notes.append(f"Key question: {key_question}")
    if rationale:
        notes.append(f"\nWhy this lens fits: {rationale}")
    agent_prompt = primary_lens.get('agent_prompt', '')
    if agent_prompt:
        notes.append(f"\nAGENT PROMPT ({ctx.key_stage}):\n{agent_prompt}")
    if question_stems:
        notes.append(f"\nQUESTION STEMS ({ctx.key_stage}):")
        for s in question_stems:
            notes.append(f"- {s}")
    if len(ctx.thinking_lenses) > 1:
        notes.append(f"\nSECONDARY LENS: {ctx.thinking_lenses[1].get('lens_name', '')}")
        notes.append(f"Key question: {ctx.thinking_lenses[1].get('key_question', '')}")
        sec_rationale = ctx.thinking_lenses[1].get('rationale', '')
        if sec_rationale:
            notes.append(f"Rationale: {sec_rationale}")
    set_notes(slide, '\n'.join(notes))


# ── Session structure slide ─────────────────────────────────────────

def _build_session_structure_slide(prs, ctx, pal, age):
    """Session structure from vehicle templates."""
    if not ctx.templates:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pal['bg'])

    emojis = _emojis(ctx.subject)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 f'{emojis[4]}  Session structure', font_size=age['title_size'],
                 colour=DARK_BROWN, bold=True)
    add_accent_bar(slide, Inches(1.1), pal['primary'])

    y = Inches(1.5)
    for vt in ctx.templates[:2]:
        name = vt.get('name', 'Template')
        phases = vt.get('session_structure', [])
        assessment = vt.get('assessment_approach', '')

        # Template name
        add_text_box(slide, Inches(0.5), y, Inches(11), Inches(0.5),
                     name, font_size=age['body_size'] + 4,
                     colour=pal['primary'], bold=True)
        y += Inches(0.6)

        # Phase arrows
        if phases:
            cols = min(len(phases), 5)
            phase_w = Inches(12 / cols - 0.3)
            for i, phase in enumerate(phases[:5]):
                left = Inches(0.3 + i * (12.5 / cols))
                colour = pal['primary'] if i % 2 == 0 else pal['accent']
                add_card(slide, left, y, phase_w, Inches(1.2), colour)
                add_text_box(slide, left + Inches(0.1), y + Inches(0.1),
                             phase_w - Inches(0.2), Inches(1.0),
                             phase, font_size=age['card_font_body'],
                             colour=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                             auto_shrink=True)
            y += Inches(1.6)

        # Assessment approach
        if assessment:
            add_text_box(slide, Inches(0.5), y, Inches(11), Inches(0.8),
                         f"Assessment: {assessment}",
                         font_size=age['body_size'] - 2, colour=DARK_SLATE,
                         auto_shrink=True)
            y += Inches(0.8)

    # Teacher notes
    notes = ['SESSION STRUCTURE\n']
    for vt in ctx.templates:
        notes.append(f"Template: {vt.get('name', '')}")
        phases = vt.get('session_structure', [])
        if phases:
            notes.append(f"Phases: {' → '.join(phases)}")
        assess = vt.get('assessment_approach', '')
        if assess:
            notes.append(f"Assessment: {assess}")
        ap = vt.get('ks_agent_prompt', '') or vt.get('agent_prompt', '')
        if ap:
            notes.append(f"\nAgent prompt ({ctx.key_stage}):\n{ap}")
        stems = vt.get('ks_question_stems', [])
        if stems:
            notes.append(f"\nTemplate question stems:")
            for s in stems:
                notes.append(f"- {s}")
        notes.append('')
    set_notes(slide, '\n'.join(notes))


# ── Graph context slide (for teachers) ──────────────────────────────

def _build_graph_context_slide(prs, ctx, pal, age):
    """Final slide: graph node IDs and query hints for teachers who want more data."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, STONE_GREY)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 '📊  Graph context (for data lookup)', font_size=age['title_size'] - 8,
                 colour=DARK_SLATE, bold=True)
    add_accent_bar(slide, Inches(1.0), DARK_SLATE, height=Inches(0.04))

    # Node info card
    add_card(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(3.5), WHITE,
             border_colour=DARK_SLATE)

    info_lines = [
        f"Node type: {ctx.label}",
        f"Study ID: {ctx.study_id}",
        f"Subject: {ctx.subject}",
        f"Key stage: {ctx.key_stage}",
        f"Concepts: {len(ctx.concepts)} ({sum(1 for c in ctx.concepts if c.get('is_primary'))} primary)",
    ]
    domain_ids = ctx.study.get('domain_ids', [])
    if isinstance(domain_ids, str):
        try:
            domain_ids = json.loads(domain_ids)
        except (json.JSONDecodeError, TypeError):
            domain_ids = []
    if domain_ids:
        info_lines.append(f"Domain IDs: {', '.join(domain_ids[:4])}")
    template_ids = ctx.study.get('uses_template', [])
    if isinstance(template_ids, str):
        template_ids = [template_ids]
    if template_ids:
        info_lines.append(f"Template IDs: {', '.join(template_ids)}")
    if ctx.thinking_lenses:
        info_lines.append(f"Thinking lens: {ctx.thinking_lenses[0].get('lens_name', '')}")

    y_info = Inches(1.5)
    for line in info_lines:
        add_text_box(slide, Inches(0.7), y_info, Inches(5.3), Inches(0.35),
                     line, font_size=14, colour=DARK_SLATE)
        y_info += Inches(0.35)

    # Cypher query hints card
    add_card(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(3.5), WHITE,
             border_colour=DARK_SLATE)

    id_field = {
        'HistoryStudy': 'study_id', 'GeoStudy': 'study_id',
        'ScienceEnquiry': 'enquiry_id', 'EnglishUnit': 'unit_id',
    }.get(ctx.label, 'suggestion_id')

    queries = [
        f"// Find this study node",
        f"MATCH (ts:{ctx.label} {{{id_field}: '{ctx.study_id}'}})",
        f"RETURN ts",
        f"",
        f"// Get concepts + difficulty levels",
        f"MATCH (ts:{ctx.label} {{{id_field}: '{ctx.study_id}'}})",
        f"  -[:DELIVERS_VIA]->(c:Concept)",
        f"  -[:HAS_DIFFICULTY_LEVEL]->(dl)",
        f"RETURN c.name, dl.label, dl.description",
    ]

    y_q = Inches(1.5)
    for line in queries:
        add_text_box(slide, Inches(7.0), y_q, Inches(5.5), Inches(0.3),
                     line, font_size=11, colour=DARK_SLATE)
        y_q += Inches(0.28)

    # Source documents
    if ctx.source_documents:
        sd = ctx.source_documents[0]
        sd_name = sd.get('name', '')
        add_text_box(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(0.4),
                     f"Source document: {sd_name}",
                     font_size=12, colour=DARK_SLATE)

    # Epistemic skills
    if ctx.epistemic_skills:
        skill_names = [s.get('name', '') for s in ctx.epistemic_skills[:5]]
        add_text_box(slide, Inches(0.5), Inches(5.7), Inches(12), Inches(0.4),
                     f"Disciplinary skills: {', '.join(skill_names)}",
                     font_size=12, colour=DARK_SLATE)

    # Teacher notes with full agent prompt context
    notes = ['GRAPH CONTEXT — Use these IDs to look up additional data in the curriculum graph.\n']
    notes.append(f"Node: {ctx.label} | ID: {ctx.study_id} | Subject: {ctx.subject} | KS: {ctx.key_stage}")
    notes.append(f"\nConcept IDs delivered:")
    for c in ctx.concepts:
        cid = c.get('concept_id', '')
        name = c.get('concept_name', c.get('name', ''))
        pri = '(PRIMARY)' if c.get('is_primary') else ''
        notes.append(f"  {cid}: {name} {pri}")
    if domain_ids:
        notes.append(f"\nDomain IDs: {', '.join(domain_ids)}")
    if ctx.source_documents:
        notes.append(f"\nSource documents:")
        for sd in ctx.source_documents:
            notes.append(f"  {sd.get('name', '')} ({sd.get('reference', '')})")
    if ctx.thinking_lenses:
        notes.append(f"\nThinking lenses:")
        for tl in ctx.thinking_lenses:
            notes.append(f"  {tl.get('lens_name', '')} (rank {tl.get('rank', '')}): {tl.get('key_question', '')}")
    if ctx.epistemic_skills:
        notes.append(f"\nEpistemic skills:")
        for s in ctx.epistemic_skills:
            notes.append(f"  {s.get('name', '')} ({s.get('skill_type', '')})")
    set_notes(slide, '\n'.join(notes))


# ── Science-specific slides ─────────────────────────────────────────

def _build_science_investigation_slide(prs, ctx, pal, age):
    """Science investigation detail slide."""
    study = ctx.study
    variables = study.get('variables')
    if isinstance(variables, str):
        try:
            variables = json.loads(variables)
        except (json.JSONDecodeError, TypeError):
            variables = None
    equipment = _as_list(study.get('equipment'))
    safety = study.get('safety_notes', '')

    if not variables and not equipment:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pal['bg'])

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 '🔬  Investigation plan', font_size=age['title_size'],
                 colour=DARK_BROWN, bold=True)
    add_accent_bar(slide, Inches(1.2), pal['primary'])

    # Variables as a structured layout
    if variables and isinstance(variables, dict):
        var_items = [
            ('Independent variable', variables.get('independent', ''), pal['primary']),
            ('Dependent variable', variables.get('dependent', ''), pal['accent']),
        ]
        controlled = variables.get('controlled', [])
        if isinstance(controlled, list):
            var_items.append(('Controlled variables', ', '.join(controlled), DARK_SLATE))
        elif controlled:
            var_items.append(('Controlled variables', str(controlled), DARK_SLATE))

        for i, (label, val, colour) in enumerate(var_items):
            left = Inches(0.5 + i * 4.2)
            add_card(slide, left, Inches(1.5), Inches(3.8), Inches(1.8), WHITE,
                     border_colour=colour)
            add_text_box(slide, left + Inches(0.1), Inches(1.6), Inches(3.6), Inches(0.4),
                         label, font_size=age['card_font_body'] - 2,
                         colour=colour, bold=True)
            add_text_box(slide, left + Inches(0.1), Inches(2.1), Inches(3.6), Inches(1.0),
                         val, font_size=age['card_font_body'],
                         colour=DARK_BROWN, auto_shrink=True)

    # Equipment list
    if equipment:
        add_text_box(slide, Inches(0.5), Inches(3.6), Inches(5), Inches(0.5),
                     '🧪 Equipment needed:', font_size=age['body_size'],
                     colour=DARK_BROWN, bold=True)
        add_bullet_list(slide, Inches(0.5), Inches(4.1), Inches(5), Inches(2.5),
                        equipment[:8], font_size=age['bullet_size'] - 2,
                        emoji_bullets=False)

    # Safety box
    if safety:
        add_card(slide, Inches(6), Inches(3.6), Inches(6.5), Inches(3),
                 RGBColor(0xFF, 0xEB, 0xEE), border_colour=RGBColor(0xE7, 0x4C, 0x3C))
        add_text_box(slide, Inches(6.2), Inches(3.7), Inches(6.1), Inches(0.4),
                     '⚠️ Safety notes', font_size=age['body_size'],
                     colour=RGBColor(0xC0, 0x39, 0x2B), bold=True)
        add_text_box(slide, Inches(6.2), Inches(4.2), Inches(6.1), Inches(2.2),
                     safety, font_size=age['body_size'] - 2,
                     colour=DARK_BROWN, auto_shrink=True)

    # Teacher notes
    notes = ['INVESTIGATION PLAN SLIDE\n']
    if variables and isinstance(variables, dict):
        notes.append(f"Independent: {variables.get('independent', '')}")
        notes.append(f"Dependent: {variables.get('dependent', '')}")
        notes.append(f"Controlled: {variables.get('controlled', '')}")
    if equipment:
        notes.append(f"\nEquipment: {', '.join(equipment)}")
    if safety:
        notes.append(f"\nSafety: {safety}")
    outcome = study.get('expected_outcome', '')
    if outcome:
        notes.append(f"\nExpected outcome: {outcome}")
    recording = _as_list(study.get('recording_format'))
    if recording:
        notes.append(f"Recording format: {', '.join(recording)}")
    # Misconceptions from reference nodes
    misconceptions = ctx.references.get('misconceptions', [])
    if misconceptions:
        notes.append('\nKNOWN MISCONCEPTIONS:')
        for m in misconceptions:
            notes.append(f"- {m.get('name', '')}: {m.get('description', '')}")
    set_notes(slide, '\n'.join(notes))


# ── Geography-specific slides ───────────────────────────────────────

def _build_geography_places_slide(prs, ctx, pal, age):
    """Geography locations and fieldwork slide."""
    places = ctx.references.get('places', [])
    contrasts = ctx.references.get('contrasts', [])
    fieldwork = ctx.study.get('fieldwork_potential', '')

    if not places and not contrasts and not fieldwork:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pal['bg'])

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 '🗺️  Locations and fieldwork', font_size=age['title_size'],
                 colour=DARK_BROWN, bold=True)
    add_accent_bar(slide, Inches(1.2), pal['primary'])

    y = Inches(1.5)
    # Places as cards
    for i, place in enumerate(places[:4]):
        left = Inches(0.5 + (i % 2) * 6.3)
        if i > 0 and i % 2 == 0:
            y += Inches(2.4)
        add_card(slide, left, y, Inches(5.8), Inches(2.0), WHITE,
                 border_colour=pal['primary'])
        add_text_box(slide, left + Inches(0.15), y + Inches(0.1),
                     Inches(5.5), Inches(0.4),
                     f"📍 {place.get('name', '')}",
                     font_size=age['card_font_title'], colour=pal['primary'], bold=True)
        desc = place.get('description', '')
        if desc:
            add_text_box(slide, left + Inches(0.15), y + Inches(0.55),
                         Inches(5.5), Inches(1.3),
                         desc[:180], font_size=age['card_font_body'] - 2,
                         colour=DARK_BROWN, auto_shrink=True)

    # Fieldwork hint
    if fieldwork:
        add_card(slide, Inches(0.5), Inches(5.8), Inches(12), Inches(1.2),
                 pal['accent'])
        add_text_box(slide, Inches(0.7), Inches(5.9), Inches(11.5), Inches(1.0),
                     f"🥾 Fieldwork potential: {fieldwork}",
                     font_size=age['body_size'], colour=WHITE, auto_shrink=True)

    notes = ['LOCATIONS AND FIELDWORK\n']
    for p in places:
        notes.append(f"Place: {p.get('name', '')} — {p.get('description', '')}")
    for c in contrasts:
        notes.append(f"Contrast: {c.get('name', '')} — {c.get('description', '')}")
    if fieldwork:
        notes.append(f"\nFieldwork: {fieldwork}")
    scale = ctx.study.get('scale', '')
    if scale:
        notes.append(f"Scale: {scale}")
    set_notes(slide, '\n'.join(notes))


# ── History-specific slides ──────────────────────────────────────────

def _build_history_sources_slide(prs, ctx, pal, age):
    """Primary sources slide (History only)."""
    sources = ctx.references.get('sources', [])
    if not sources:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pal['bg'])

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 '📜  Primary sources', font_size=age['title_size'],
                 colour=DARK_BROWN, bold=True)
    add_accent_bar(slide, Inches(1.2), RGBColor(0x2E, 0x86, 0xC1))

    for i, src in enumerate(sources[:3]):
        left = Inches(0.3 + i * 4.3)
        name = src.get('name', src.get('source_name', ''))
        stype = src.get('source_type', '').replace('_', ' ')
        provenance = src.get('provenance', src.get('description', ''))

        card_colour = [RGBColor(0xFE, 0xF9, 0xE7), pal['bg'],
                       RGBColor(0xF2, 0xF3, 0xF4)][i % 3]
        add_card(slide, left, Inches(1.5), Inches(4), Inches(5.2), card_colour,
                 border_colour=pal['accent'])

        add_text_box(slide, left + Inches(0.15), Inches(1.6), Inches(3.7), Inches(0.6),
                     name, font_size=age['card_font_title'], colour=pal['primary'], bold=True)
        add_text_box(slide, left + Inches(0.15), Inches(2.2), Inches(3.7), Inches(0.4),
                     stype, font_size=age['card_font_body'] - 2,
                     colour=RGBColor(0x2E, 0x86, 0xC1))
        if provenance:
            add_text_box(slide, left + Inches(0.15), Inches(2.7), Inches(3.7), Inches(3.5),
                         provenance[:250], font_size=age['card_font_body'] - 2,
                         colour=DARK_BROWN)

    notes = ['PRIMARY SOURCES SLIDE\n']
    for src in sources:
        notes.append(f"Source: {src.get('name', '')}")
        notes.append(f"  Type: {src.get('source_type', '')}")
        ped = src.get('pedagogical_use', '')
        if ped:
            notes.append(f"  How to use: {ped}")
        loc = src.get('location', '')
        if loc:
            notes.append(f"  Location: {loc}")
        url = src.get('url', '')
        if url:
            notes.append(f"  URL: {url}")
        notes.append('')
    set_notes(slide, '\n'.join(notes))


# ── Prior knowledge slide (teacher notes) ────────────────────────────

def _build_prior_knowledge_slide(prs, ctx, pal, age):
    """Prior knowledge slide — prerequisites grouped by target concept.

    Teacher-facing: listed as teacher notes on a simple 'Prior knowledge'
    heading slide so it doesn't clutter child-facing content.
    """
    if not ctx.prerequisites:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pal['bg'])

    emoji = _emojis(ctx.subject)[6]
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 f'{emoji}  What you already know', font_size=age['title_size'],
                 colour=DARK_BROWN, bold=True)
    add_accent_bar(slide, Inches(1.2), pal['primary'])

    # Show unique prerequisite names on the slide (child-friendly)
    seen = set()
    prereq_names = []
    for pr in ctx.prerequisites:
        name = pr.get('prereq_name', '')
        if name and name not in seen:
            seen.add(name)
            prereq_names.append(name)

    add_bullet_list(slide, Inches(0.7), Inches(1.5), Inches(11), Inches(5),
                    prereq_names[:age['max_bullet_items']],
                    font_size=age['bullet_size'],
                    emoji_bullets=age['emoji_bullets'],
                    emojis=_emojis(ctx.subject))

    # Teacher notes — grouped by target concept
    notes = ['PRIOR KNOWLEDGE — What pupils should already know\n']
    grouped = {}
    for pr in ctx.prerequisites:
        target = pr.get('target_name', 'Unknown')
        grouped.setdefault(target, []).append(pr)
    for target, items in grouped.items():
        notes.append(f"\nFor concept: {target}")
        for pr in items:
            pid = pr.get('prereq_id', '')
            pname = pr.get('prereq_name', '')
            pdesc = pr.get('prereq_description', '')
            line = f"  - {pname} ({pid})"
            if pdesc:
                line += f": {pdesc[:120]}"
            notes.append(line)
    set_notes(slide, '\n'.join(notes))


# ── Assessment codes (KS2 teacher notes) ────────────────────────────

def _add_assessment_codes_to_notes(notes_lines, ctx):
    """Append KS2 ContentDomainCode info to an existing notes list."""
    if not ctx.assessment_codes:
        return
    notes_lines.append('\nKS2 ASSESSMENT CODES:')
    for ac in ctx.assessment_codes:
        code_id = ac.get('code_id', '')
        desc = ac.get('description', '')
        concept = ac.get('concept_name', '')
        line = f"  KS2 Assessment: {code_id} — {desc}"
        if concept:
            line += f" (concept: {concept})"
        notes_lines.append(line)


# ── Knowledge organiser slide ────────────────────────────────────────

def _build_knowledge_organiser_slide(prs, ctx, pal, age):
    """Knowledge organiser — compact child-facing reference slide.

    Shows definitions, key events, key figures, and period if available.
    Only renders if there is enough data to justify the slide.
    """
    study = ctx.study
    definitions = _as_list(study.get('definitions'))
    key_events = _as_list(study.get('key_events'))
    key_figures = _as_list(study.get('key_figures'))
    period = study.get('period', '')

    # Only build if we have at least two categories of content
    sections = sum([
        len(definitions) > 0,
        len(key_events) > 0,
        len(key_figures) > 0,
        bool(period),
    ])
    if sections < 2:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pal['bg'])

    emoji = _emojis(ctx.subject)[7]
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 f'{emoji}  Knowledge organiser', font_size=age['title_size'],
                 colour=DARK_BROWN, bold=True)
    add_accent_bar(slide, Inches(1.0), pal['primary'])

    y = Inches(1.3)

    # Period / timeline header
    if period:
        add_card(slide, Inches(0.5), y, Inches(12), Inches(0.7), pal['primary'])
        add_text_box(slide, Inches(0.7), y + Inches(0.1), Inches(11.5), Inches(0.5),
                     f"Period: {period}", font_size=age['body_size'],
                     colour=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        y += Inches(0.9)

    # Key terms grid (left column)
    col_left = Inches(0.5)
    col_right = Inches(6.7)
    col_w = Inches(5.8)

    if definitions:
        add_text_box(slide, col_left, y, col_w, Inches(0.4),
                     'Key terms', font_size=age['body_size'] - 2,
                     colour=pal['primary'], bold=True)
        term_y = y + Inches(0.45)
        for item in definitions[:6]:
            if isinstance(item, dict):
                word = item.get('term', item.get('word', ''))
                defn = item.get('meaning', item.get('definition', ''))
            else:
                word = str(item)
                defn = ''
            text = f"{word}: {defn}" if defn else word
            if len(text) > 80:
                text = text[:77] + '...'
            add_text_box(slide, col_left + Inches(0.1), term_y, col_w - Inches(0.2),
                         Inches(0.35), f"• {text}",
                         font_size=age['card_font_body'] - 2, colour=DARK_BROWN)
            term_y += Inches(0.35)

    # Key events (right column or below if no definitions)
    events_col = col_right if definitions else col_left
    if key_events:
        add_text_box(slide, events_col, y, col_w, Inches(0.4),
                     'Key events', font_size=age['body_size'] - 2,
                     colour=pal['accent'], bold=True)
        ev_y = y + Inches(0.45)
        for event in key_events[:6]:
            if len(event) > 80:
                event = event[:77] + '...'
            add_text_box(slide, events_col + Inches(0.1), ev_y, col_w - Inches(0.2),
                         Inches(0.35), f"• {event}",
                         font_size=age['card_font_body'] - 2, colour=DARK_BROWN)
            ev_y += Inches(0.35)

    # Key figures at bottom
    if key_figures:
        fig_y = Inches(5.2)
        add_text_box(slide, Inches(0.5), fig_y, Inches(12), Inches(0.4),
                     'Key figures', font_size=age['body_size'] - 2,
                     colour=pal['primary'], bold=True)
        fig_y += Inches(0.45)
        cols = min(len(key_figures), 4)
        for i, fig in enumerate(key_figures[:4]):
            left = Inches(0.5 + i * (12 / cols))
            card_w = Inches(12 / cols - 0.3)
            add_card(slide, left, fig_y, card_w, Inches(0.7), WHITE,
                     border_colour=pal['primary'])
            if len(fig) > 40:
                fig = fig[:37] + '...'
            add_text_box(slide, left + Inches(0.1), fig_y + Inches(0.1),
                         card_w - Inches(0.2), Inches(0.5),
                         f"👤 {fig}", font_size=age['card_font_body'] - 2,
                         colour=DARK_BROWN, bold=True)

    notes = ['KNOWLEDGE ORGANISER — Print for exercise books / display wall.']
    if period:
        notes.append(f"Period: {period}")
    if definitions:
        terms = [d.get('term', d.get('word', str(d))) if isinstance(d, dict) else str(d)
                 for d in definitions]
        notes.append(f"Key terms: {', '.join(terms)}")
    if key_events:
        notes.append(f"Key events: {', '.join(key_events)}")
    if key_figures:
        notes.append(f"Key figures: {', '.join(key_figures)}")
    set_notes(slide, '\n'.join(notes))


# ── Main render function ─────────────────────────────────────────────

def render_pptx(ctx: StudyContext) -> Presentation:
    """Render a complete PPTX presentation for a study. Returns Presentation object."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    pal = _get_palette(ctx.subject)
    age = _get_age(ctx.key_stage)

    # Universal slides
    _build_title_slide(prs, ctx, pal, age)
    _build_knowledge_organiser_slide(prs, ctx, pal, age)
    _build_prior_knowledge_slide(prs, ctx, pal, age)
    _build_concepts_slide(prs, ctx, pal, age)
    _build_enquiry_slide(prs, ctx, pal, age)
    _build_subject_content_slide(prs, ctx, pal, age)

    # Subject-specific slides
    if ctx.subject == 'History':
        _build_history_sources_slide(prs, ctx, pal, age)
    elif ctx.subject == 'Science':
        _build_science_investigation_slide(prs, ctx, pal, age)
    elif ctx.subject == 'Geography':
        _build_geography_places_slide(prs, ctx, pal, age)

    # Thinking lens (cognitive framing)
    _build_thinking_lens_slide(prs, ctx, pal, age)

    # Session structure (vehicle templates)
    _build_session_structure_slide(prs, ctx, pal, age)

    # Universal closing slides
    _build_vocabulary_slide(prs, ctx, pal, age)
    _build_assessment_slide(prs, ctx, pal, age)
    _build_next_slide(prs, ctx, pal, age)

    # Graph context (teacher reference)
    _build_graph_context_slide(prs, ctx, pal, age)

    return prs
