"""
render.py — Build PPTX from story.json + curriculum_unit.json

The story brief determines visual treatment.
The curriculum unit provides the content.
No intermediate deck.json — straight to PPTX.

Usage:
    python3 render.py unit3-protectionism/curriculum_unit.json \
                      unit3-protectionism/story.json \
                      [output.pptx]
"""

import json
import sys
from io import BytesIO
from pathlib import Path
from html import escape as xe
from lxml import etree
from pptx import Presentation
from pptx.util import Emu

# ── EMU helpers ───────────────────────────────────────────────────────────────

def emu(inches): return int(inches * 914400)

W = 12192000   # exact standard widescreen 16:9 (13.333... inches × 914400)
H = 6858000    # exact standard widescreen 16:9 (7.5 inches × 914400)

# ── Palette ───────────────────────────────────────────────────────────────────

NAVY   = "121F47"
TEAL   = "009FA3"
SLATE  = "2B576B"
WHITE  = "FFFFFF"
GREY   = "656565"
L_GREY = "F4F4F4"
AMBER  = "F59E0B"
RED    = "DC2626"
GREEN  = "059669"
L_TEAL = "E0F7F7"
YELLOW = "FFFBEB"

# Beat → background colour
BEAT_BG = {
    "hook":             NAVY,
    "moment":           NAVY,
    "data_reveal":      NAVY,
    "pivot":            NAVY,
    "threat":           "1A0F00",   # near-black amber tint
    "twist":            SLATE,
}

# ── Layout constants ──────────────────────────────────────────────────────────

ML = emu(0.50)
MR = emu(0.50)
CW = W - ML - MR
MT = emu(0.20)

LOGO_W = emu(1.7)
LOGO_H = int(LOGO_W * 108 / 385)   # harris-logo.png actual ratio: 385×108
LOGO_X = W - LOGO_W - emu(0.30)
LOGO_Y = H - LOGO_H - emu(0.15)

# ── ID counter ────────────────────────────────────────────────────────────────

_id = [2]

def nid():
    _id[0] += 1
    return _id[0]

def reset_ids():
    _id[0] = 2

# ── XML primitives ────────────────────────────────────────────────────────────

def solid(h): return f'<a:solidFill><a:srgbClr val="{h}"/></a:solidFill>'
def nofill():  return '<a:noFill/>'
def noline():  return '<a:ln><a:noFill/></a:ln>'

def line(h, w=12700):
    return f'<a:ln w="{w}"><a:solidFill><a:srgbClr val="{h}"/></a:solidFill></a:ln>'

def run(text, pt, bold=False, italic=False, color=WHITE, face="Calibri"):
    b = "1" if bold else "0"
    i = "1" if italic else "0"
    sz = int(pt * 100)
    fill = solid(color) if color else ""
    return (
        f'<a:r><a:rPr lang="en-GB" sz="{sz}" b="{b}" i="{i}" dirty="0">'
        f'<a:latin typeface="{face}"/>{fill}</a:rPr>'
        f'<a:t>{xe(str(text))}</a:t></a:r>'
    )

def runL(text, pt, bold=False, italic=False, color=WHITE):
    return run(text, pt, bold, italic, color, face="Calibri Light")

def para(runs_xml, align="l", spc_bef=0, spc_aft=0, lnspc=115):
    alg = {"l":"l","c":"ctr","r":"r","left":"l","center":"ctr"}.get(align,"l")
    bef = f'<a:spcBef><a:spcPts val="{spc_bef}"/></a:spcBef>' if spc_bef else ""
    aft = f'<a:spcAft><a:spcPts val="{spc_aft}"/></a:spcAft>' if spc_aft else ""
    ln  = f'<a:lnSpc><a:spcPct val="{lnspc*1000}"/></a:lnSpc>'
    return (f'<a:p><a:pPr algn="{alg}">{bef}{aft}{ln}<a:buNone/></a:pPr>'
            f'{runs_xml}</a:p>')

def bul(runs_xml, level=1, spc_bef=0, spc_aft=80):
    margin = 304800 if level == 1 else 609600
    indent = -304800
    char   = "▪" if level == 1 else "–"
    bef = f'<a:spcBef><a:spcPts val="{spc_bef}"/></a:spcBef>' if spc_bef else ""
    aft = f'<a:spcAft><a:spcPts val="{spc_aft}"/></a:spcAft>'
    ln  = '<a:lnSpc><a:spcPct val="115000"/></a:lnSpc>'
    return (f'<a:p><a:pPr marL="{margin}" indent="{indent}">{bef}{aft}{ln}'
            f'<a:buFont typeface="Calibri"/><a:buChar char="{char}"/></a:pPr>'
            f'{runs_xml}</a:p>')

def empty_para():
    return '<a:p><a:pPr><a:buNone/></a:pPr><a:endParaRPr lang="en-GB" dirty="0"/></a:p>'

def txBody(paras, anchor="t", il=91440, it=45720, ir=91440, ib=45720):
    return (f'<p:txBody>'
            f'<a:bodyPr wrap="square" lIns="{il}" tIns="{it}" rIns="{ir}" bIns="{ib}" anchor="{anchor}" anchorCtr="0">'
            f'<a:normAutofit/></a:bodyPr><a:lstStyle/>{paras}</p:txBody>')

def rect(x, y, w, h, fill=WHITE, ln_h=None, name="rect"):
    sid = nid()
    fill_xml = solid(fill)
    ln_xml   = line(ln_h) if ln_h else noline()
    return (f'<p:sp><p:nvSpPr><p:cNvPr id="{sid}" name="{name}"/>'
            f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>'
            f'<p:spPr><a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{w}" cy="{h}"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>{fill_xml}{ln_xml}</p:spPr>'
            f'<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>')

def txt(x, y, w, h, tb, fill=None, name="txt"):
    sid = nid()
    fill_xml = solid(fill) if fill else nofill()
    return (f'<p:sp><p:nvSpPr><p:cNvPr id="{sid}" name="{name}"/>'
            f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>'
            f'<p:spPr><a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{w}" cy="{h}"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>{fill_xml}{noline()}</p:spPr>'
            f'{tb}</p:sp>')

def pic(x, y, w, h, rel_id, name="img"):
    """Place image at exact dimensions. Caller is responsible for matching the image's
    aspect ratio — use fit_image_box() to calculate correct w/h before calling this."""
    sid = nid()
    return (f'<p:pic><p:nvPicPr><p:cNvPr id="{sid}" name="{name}"/>'
            f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr><p:nvPr/></p:nvPicPr>'
            f'<p:blipFill><a:blip r:embed="{rel_id}"/>'
            f'<a:stretch><a:fillRect/></a:stretch></p:blipFill>'
            f'<p:spPr><a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{w}" cy="{h}"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>')


def fit_image_box(slot_x, slot_y, slot_w, slot_h, img_ratio_w, img_ratio_h, align="center"):
    """
    Given a slot (max area) and an image's aspect ratio, return (x, y, w, h) for the
    image box that fits inside the slot without distortion — centred by default.

    img_ratio_w / img_ratio_h define the ratio (e.g. 16, 9 for 16:9).

    Usage:
        # For a 16:9 image in a half-slide slot:
        ix, iy, iw, ih = fit_image_box(ML, top, slot_w, slot_h, 16, 9)
        shapes += pic(ix, iy, iw, ih, rel_id)

    When requesting images from an MCP, compute pixel dimensions from the output:
        px_w = int(iw / 914400 * 96)   # 96 dpi
        px_h = int(ih / 914400 * 96)
        # Then request image at px_w × px_h (or equivalent --ar ratio)
    """
    ratio = img_ratio_w / img_ratio_h
    slot_ratio = slot_w / slot_h

    if ratio > slot_ratio:
        # Image is wider than slot — constrain by width
        fit_w = slot_w
        fit_h = int(slot_w / ratio)
    else:
        # Image is taller than slot — constrain by height
        fit_h = slot_h
        fit_w = int(slot_h * ratio)

    if align == "center":
        fit_x = slot_x + (slot_w - fit_w) // 2
        fit_y = slot_y + (slot_h - fit_h) // 2
    elif align == "left":
        fit_x = slot_x
        fit_y = slot_y + (slot_h - fit_h) // 2
    elif align == "top":
        fit_x = slot_x + (slot_w - fit_w) // 2
        fit_y = slot_y
    else:
        fit_x, fit_y = slot_x, slot_y

    return fit_x, fit_y, fit_w, fit_h


def image_request_dimensions(slot_w_emu, slot_h_emu, img_ratio_w, img_ratio_h, dpi=150):
    """
    Given a slot in EMU and a desired image ratio, return (px_w, px_h) to request
    from an image MCP (Midjourney --ar or edu-images size parameter).

    Example:
        px_w, px_h = image_request_dimensions(slot_w, slot_h, 16, 9)
        # Request: Midjourney  --ar 16:9  or  width=px_w height=px_h
    """
    _, _, fit_w, fit_h = fit_image_box(0, 0, slot_w_emu, slot_h_emu, img_ratio_w, img_ratio_h)
    px_w = max(256, int(fit_w / 914400 * dpi))
    px_h = max(256, int(fit_h / 914400 * dpi))
    return px_w, px_h

def slide_wrap(bg, shapes, notes=""):
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        '  xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        '  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        '<p:cSld>'
        f'<p:bg><p:bgPr>{solid(bg)}<a:effectLst/></p:bgPr></p:bg>'
        '<p:spTree>'
        '<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
        f'<p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="{W}" cy="{H}"/>'
        f'<a:chOff x="0" y="0"/><a:chExt cx="{W}" cy="{H}"/></a:xfrm></p:grpSpPr>'
        f'{shapes}'
        '</p:spTree></p:cSld>'
        '<p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>'
        '</p:sld>'
    )

# ── Shared components ─────────────────────────────────────────────────────────

def logo_shape(rel_id):
    if not rel_id:
        return ""
    return pic(LOGO_X, LOGO_Y, LOGO_W, LOGO_H, rel_id, name="logo")

def header_bar(title_text, bg=NAVY, fg=WHITE, subtitle=None):
    """Full-width header bar with title (and optional subtitle)."""
    bar_h = emu(1.1) if not subtitle else emu(1.4)
    shapes = rect(0, 0, W, bar_h, fill=bg, name="header_bg")
    # Teal left accent
    shapes += rect(0, 0, emu(0.12), bar_h, fill=TEAL, name="accent")
    title_p = para(runL(title_text, 24, bold=True, color=fg), align="l")
    shapes += txt(ML, emu(0.12) if not subtitle else emu(0.08),
                  CW, emu(0.62), txBody(title_p, anchor="ctr"), name="title")
    if subtitle:
        sub_p = para(run(subtitle, 11, color=TEAL, face="Calibri"), align="l")
        shapes += txt(ML, emu(0.72), CW, emu(0.45), txBody(sub_p, anchor="t"), name="subtitle")
    return shapes, bar_h

def callout_box(text, y, color=L_TEAL, text_color=NAVY):
    h = emu(0.85)
    shapes  = rect(ML, y, CW, h, fill=color, name="callout_bg")
    shapes += rect(ML, y, emu(0.06), h, fill=TEAL, name="callout_bar")
    p = para(run(text, 11, italic=True, color=text_color), align="l", spc_bef=0)
    shapes += txt(ML + emu(0.15), y + emu(0.10), CW - emu(0.20), h - emu(0.20),
                  txBody(p, anchor="ctr"), name="callout")
    return shapes

def bullets_block(bullets, x, y, w, h, base_sz=13, color=NAVY):
    paras = ""
    for b in bullets:
        lvl   = b.get("level", 1)
        text  = b.get("text", "")
        hi    = b.get("highlight", False)
        sz    = base_sz - (1 if lvl > 1 else 0)
        col   = TEAL if hi else color
        bold  = hi
        paras += bul(run(text, sz, bold=bold, color=col), level=lvl, spc_aft=60)
    paras = paras or empty_para()
    return txt(x, y, w, h, txBody(paras, anchor="t"), name="bullets")

def ao_tag_badge(ao, x, y):
    """Small AO badge — AO1/AO2/AO3/AO4"""
    colors = {"AO1": SLATE, "AO2": TEAL, "AO3": GREEN, "AO4": AMBER}
    col = colors.get(ao, GREY)
    w, h = emu(0.55), emu(0.28)
    shapes  = rect(x, y, w, h, fill=col, name=f"badge_{ao}")
    p = para(run(ao, 8, bold=True, color=WHITE), align="c")
    shapes += txt(x, y, w, h, txBody(p, anchor="ctr"), name=f"badge_txt_{ao}")
    return shapes

# ── Beat-type renderers ───────────────────────────────────────────────────────

def render_hook(brief, sl, logo_rel, img_rels):
    """Hook / Moment / Pivot — dark, claim-dominant, minimal text."""
    bg = BEAT_BG.get(brief["beat_type"], NAVY)
    claim = brief.get("claim", sl.get("title",""))
    anchor = brief.get("anchor", sl.get("subtitle",""))
    energy = brief.get("energy","high")

    shapes = ""
    # Teal left rail
    shapes += rect(0, 0, emu(0.18), H, fill=TEAL, name="rail")
    # Claim — large, white, centred vertically
    claim_sz = 36 if len(claim) < 80 else 28
    claim_p = para(runL(claim, claim_sz, bold=True, color=WHITE), align="l",
                   spc_bef=0, spc_aft=200, lnspc=120)
    claim_y = emu(1.8) if anchor else emu(2.2)
    shapes += txt(ML + emu(0.18), claim_y, CW - emu(0.18), emu(3.0),
                  txBody(claim_p, anchor="ctr"), name="claim")
    # Anchor — teal, smaller
    if anchor:
        anc_p = para(run(anchor, 14, color=TEAL), align="l", spc_bef=0)
        shapes += txt(ML + emu(0.18), claim_y + emu(2.5), CW - emu(0.18), emu(1.0),
                      txBody(anc_p, anchor="t"), name="anchor")
    # Bottom line
    shapes += rect(0, H - emu(0.08), W, emu(0.08), fill=TEAL, name="btm_line")
    shapes += logo_shape(logo_rel)
    return slide_wrap(bg, shapes)


def render_data_reveal(brief, sl, logo_rel, img_rels):
    """Data reveal — one massive number, minimal context."""
    bg = NAVY
    stat  = sl.get("stat") or brief.get("claim","")
    quote = sl.get("quote") or ""
    callout = sl.get("callout","")

    shapes = ""
    # Top teal bar
    shapes += rect(0, 0, W, emu(0.12), fill=TEAL, name="top_bar")
    # Slide title
    title_p = para(runL(sl.get("title",""), 16, bold=True, color=TEAL), align="l")
    shapes += txt(ML, emu(0.18), CW, emu(0.55), txBody(title_p, anchor="ctr"), name="title")
    # Big stat
    stat_p = para(runL(stat, 54, bold=True, color=WHITE), align="c")
    shapes += txt(ML, emu(1.0), CW, emu(2.0), txBody(stat_p, anchor="ctr"), name="stat")
    # Quote
    if quote:
        q_p = para(run(quote, 13, italic=True, color=L_TEAL), align="c", lnspc=130)
        shapes += txt(ML, emu(3.3), CW, emu(1.6), txBody(q_p, anchor="t"), name="quote")
    # Callout
    if callout:
        shapes += callout_box(callout, H - emu(1.15), color=SLATE, text_color=WHITE)
    shapes += logo_shape(logo_rel)
    return slide_wrap(bg, shapes)


def render_section_header(brief, sl, logo_rel, img_rels):
    """Section opener — full navy, large section title."""
    title  = sl.get("section_title") or sl.get("title","")
    sub    = sl.get("subtitle") or brief.get("claim","")
    shapes = ""
    # Full navy bg already in slide_wrap
    # Vertical teal stripe
    shapes += rect(0, 0, emu(0.35), H, fill=TEAL, name="stripe")
    # Section label
    lbl_p = para(run("SECTION", 10, bold=True, color=TEAL), align="l")
    shapes += txt(ML + emu(0.35), emu(2.0), CW - emu(0.35), emu(0.5),
                  txBody(lbl_p, anchor="b"), name="section_label")
    # Title
    t_sz = 40 if len(title) < 35 else 32
    t_p = para(runL(title, t_sz, bold=True, color=WHITE), align="l", lnspc=115)
    shapes += txt(ML + emu(0.35), emu(2.5), CW - emu(0.35), emu(2.2),
                  txBody(t_p, anchor="t"), name="title")
    # Subtitle
    if sub:
        s_p = para(run(sub, 13, color=L_TEAL), align="l", lnspc=130)
        shapes += txt(ML + emu(0.35), emu(5.0), CW - emu(0.35), emu(1.5),
                      txBody(s_p, anchor="t"), name="subtitle")
    # Bottom teal bar
    shapes += rect(0, H - emu(0.12), W, emu(0.12), fill=TEAL, name="btm")
    shapes += logo_shape(logo_rel)
    return slide_wrap(NAVY, shapes)


def render_content(brief, sl, logo_rel, img_rels):
    """Standard content slide — header + bullets + optional callout."""
    title   = sl.get("title","")
    bullets = sl.get("bullets",[])
    callout = sl.get("callout","")
    bg      = WHITE

    shapes = ""
    hdr_shapes, bar_h = header_bar(title, bg=NAVY)
    shapes += hdr_shapes

    has_callout = bool(callout)
    callout_h   = emu(0.95) if has_callout else 0
    body_top    = bar_h + emu(0.18)
    body_h      = H - body_top - emu(0.15) - callout_h - emu(0.65)

    if bullets:
        shapes += bullets_block(bullets, ML, body_top, CW, body_h, color=NAVY)

    if has_callout:
        shapes += callout_box(callout, H - callout_h - emu(0.60))

    shapes += logo_shape(logo_rel)
    return slide_wrap(bg, shapes)


def render_two_column(brief, sl, logo_rel, img_rels):
    """Two-column tension / contrast slide."""
    title = sl.get("title","")
    left  = sl.get("left_column",{})
    right = sl.get("right_column",{})
    bg    = WHITE

    shapes = ""
    hdr_shapes, bar_h = header_bar(title, bg=NAVY)
    shapes += hdr_shapes

    top = bar_h + emu(0.20)
    col_w = (CW - emu(0.25)) / 2
    # Left column background
    shapes += rect(ML, top, int(col_w), H - top - emu(0.65), fill=L_GREY, name="left_bg")
    # Right column background
    shapes += rect(ML + int(col_w) + emu(0.25), top, int(col_w),
                   H - top - emu(0.65), fill=L_TEAL, name="right_bg")

    def col_content(col_data, x, col_color):
        if not col_data:
            return ""
        out = ""
        hdr = col_data.get("header","")
        if hdr:
            h_p = para(run(hdr, 12, bold=True, color=WHITE), align="c")
            out += txt(x, top, int(col_w), emu(0.45),
                       txBody(h_p, anchor="ctr"), fill=col_color, name="col_hdr")
        items = col_data.get("items",[])
        if items:
            paras = "".join(
                bul(run(it if isinstance(it,str) else it.get("text",""), 12, color=NAVY),
                    level=1, spc_aft=80)
                for it in items
            )
            out += txt(x + emu(0.12), top + emu(0.52), int(col_w) - emu(0.24),
                       H - top - emu(1.20), txBody(paras, anchor="t"), name="col_body")
        return out

    shapes += col_content(left,  ML,                              NAVY)
    shapes += col_content(right, ML + int(col_w) + emu(0.25),    TEAL)
    shapes += logo_shape(logo_rel)
    return slide_wrap(bg, shapes)


def render_diagram(brief, sl, logo_rel, img_rels):
    """Diagram slide — image left, annotated steps right."""
    title      = sl.get("title","")
    bullets    = sl.get("bullets",[])
    diag_img   = sl.get("diagram_image") or sl.get("image","")
    callout    = sl.get("callout","")
    bg         = WHITE

    shapes = ""
    hdr_shapes, bar_h = header_bar(title, bg=NAVY)
    shapes += hdr_shapes

    top = bar_h + emu(0.15)
    slot_h = H - top - emu(0.65)

    if diag_img and diag_img in img_rels:
        # Slot for image: left ~55% of content width
        slot_w = int(CW * 0.55)
        # All current images are 16:9 (2752×1536). fit_image_box ensures the
        # placed box matches the image ratio exactly — no distortion.
        ix, iy, iw, ih = fit_image_box(ML, top, slot_w, slot_h, 16, 9, align="top")
        shapes += pic(ix, iy, iw, ih, img_rels[diag_img], name="diagram")
        ann_x = ML + slot_w + emu(0.20)
        ann_w = CW - slot_w - emu(0.20)
    else:
        ann_x = ML
        ann_w = CW

    if bullets:
        paras = ""
        for i, b in enumerate(bullets, 1):
            txt_b = b.get("text","") if isinstance(b,dict) else str(b)
            paras += bul(run(f"{i}.  {txt_b}", 12, color=NAVY), level=1, spc_aft=80)
        shapes += txt(ann_x, top, ann_w, slot_h - emu(0.10),
                      txBody(paras, anchor="t"), name="steps")

    if callout:
        shapes += callout_box(callout, H - emu(1.05), color=YELLOW, text_color=NAVY)

    shapes += logo_shape(logo_rel)
    return slide_wrap(bg, shapes)


def render_examiner_tip(brief, sl, logo_rel, img_rels):
    """Examiner tip — amber accent, warning tone."""
    title   = sl.get("title","")
    bullets = sl.get("bullets",[])
    callout = sl.get("callout","")
    source  = sl.get("source","")
    bg      = WHITE

    shapes = ""
    # Amber left bar
    shapes += rect(0, 0, emu(0.18), H, fill=AMBER, name="amber_bar")
    # Navy header
    shapes += rect(emu(0.18), 0, W - emu(0.18), emu(1.1), fill=NAVY, name="hdr_bg")
    # "EXAMINER TIP" label
    et_p = para(run("⚠  EXAMINER TIP", 10, bold=True, color=AMBER), align="l")
    shapes += txt(ML + emu(0.18), emu(0.08), CW, emu(0.36),
                  txBody(et_p, anchor="ctr"), name="et_label")
    # Title
    t_p = para(runL(title, 22, bold=True, color=WHITE), align="l", lnspc=115)
    shapes += txt(ML + emu(0.18), emu(0.42), CW - emu(0.18), emu(0.62),
                  txBody(t_p, anchor="ctr"), name="title")

    top = emu(1.18)
    callout_h = emu(0.95) if callout else 0
    body_h = H - top - emu(0.15) - callout_h - emu(0.65)

    if bullets:
        shapes += bullets_block(bullets, ML, top, CW, body_h, color=NAVY)

    if source:
        src_p = para(run(f"Source: {source}", 9, italic=True, color=GREY), align="l")
        shapes += txt(ML, H - emu(1.50), CW, emu(0.30), txBody(src_p, anchor="t"), name="source")

    if callout:
        shapes += callout_box(callout, H - callout_h - emu(0.60), color=YELLOW, text_color=NAVY)

    shapes += logo_shape(logo_rel)
    return slide_wrap(bg, shapes)


def render_question(brief, sl, logo_rel, img_rels):
    """Question / practice slide."""
    title   = sl.get("title","")
    marks   = sl.get("marks") or ""
    q_text  = sl.get("question_text","")
    stimuli = sl.get("stimulus","")
    time_g  = sl.get("time_guidance","")
    bg      = L_GREY

    shapes = ""
    # Navy top bar
    shapes += rect(0, 0, W, emu(0.90), fill=NAVY, name="top_bar")
    # Marks badge
    if marks:
        mk_p = para(run(f"{marks} marks", 14, bold=True, color=WHITE), align="c")
        shapes += txt(W - emu(2.0), emu(0.08), emu(1.80), emu(0.75),
                      txBody(mk_p, anchor="ctr"), fill=TEAL, name="marks_badge")
    # Title
    t_p = para(runL(title, 20, bold=True, color=WHITE), align="l")
    shapes += txt(ML, emu(0.10), W - emu(2.20), emu(0.72),
                  txBody(t_p, anchor="ctr"), name="title")

    top = emu(1.0)
    if stimuli:
        s_p = para(run(stimuli, 12, italic=True, color=SLATE), align="l", lnspc=130)
        s_h = emu(1.2)
        shapes += txt(ML, top, CW, s_h, txBody(s_p, anchor="t"),
                      fill=WHITE, name="stimulus")
        top += s_h + emu(0.20)

    if q_text:
        q_p = para(run(q_text, 14, bold=True, color=NAVY), align="l", lnspc=130)
        q_h = H - top - emu(1.0)
        shapes += txt(ML, top, CW, q_h, txBody(q_p, anchor="t"), name="question")

    if time_g:
        tm_p = para(run(f"⏱  {time_g}", 11, italic=True, color=GREY), align="r")
        shapes += txt(ML, H - emu(0.80), CW, emu(0.40), txBody(tm_p, anchor="ctr"), name="timing")

    shapes += logo_shape(logo_rel)
    return slide_wrap(bg, shapes)


def render_worked_example(brief, sl, logo_rel, img_rels):
    """Worked example — step-by-step on light background."""
    title  = sl.get("title","")
    steps  = sl.get("steps",[])
    table  = sl.get("table",{})
    bg     = WHITE

    shapes = ""
    hdr_shapes, bar_h = header_bar(title, bg=NAVY)
    shapes += hdr_shapes

    top = bar_h + emu(0.20)

    # If steps, render as numbered list
    if steps:
        paras = "".join(
            bul(run(f"{i}. {s}", 13, color=NAVY), level=1, spc_aft=100)
            for i, s in enumerate(steps, 1)
        )
        shapes += txt(ML, top, CW, H - top - emu(0.65), txBody(paras, anchor="t"), name="steps")
    elif table:
        # Table as stacked rows
        rows = table.get("rows",[])
        hdr  = table.get("headers",[])
        all_rows = ([hdr] if hdr else []) + rows
        if all_rows:
            row_h = min(emu(0.55), int((H - top - emu(0.65)) / max(len(all_rows),1)))
            cols  = len(all_rows[0]) if all_rows else 1
            col_w = CW // cols
            for ri, row in enumerate(all_rows):
                is_hdr = ri == 0 and hdr
                for ci, cell in enumerate(row):
                    bg_c = NAVY if is_hdr else (L_GREY if ri % 2 == 0 else WHITE)
                    fg_c = WHITE if is_hdr else NAVY
                    cx = ML + ci * col_w
                    cy = top + ri * row_h
                    shapes += rect(cx, cy, col_w - emu(0.04), row_h, fill=bg_c, name=f"cell_{ri}_{ci}")
                    c_p = para(run(str(cell), 11, bold=is_hdr, color=fg_c), align="l")
                    shapes += txt(cx + emu(0.08), cy, col_w - emu(0.12), row_h,
                                  txBody(c_p, anchor="ctr"), name=f"cell_txt_{ri}_{ci}")

    shapes += logo_shape(logo_rel)
    return slide_wrap(bg, shapes)


def render_summary(brief, sl, logo_rel, img_rels):
    """Summary / consolidation — clean, checklist style."""
    title   = sl.get("title","")
    bullets = sl.get("bullets",[])
    bg      = WHITE

    shapes = ""
    hdr_shapes, bar_h = header_bar(title, bg=TEAL)
    shapes += hdr_shapes

    top = bar_h + emu(0.20)
    if bullets:
        paras = ""
        for b in bullets:
            text = b.get("text","") if isinstance(b,dict) else str(b)
            paras += bul(run(text, 13, color=NAVY), level=1, spc_aft=100)
        shapes += txt(ML, top, CW, H - top - emu(0.65), txBody(paras, anchor="t"), name="bullets")

    shapes += logo_shape(logo_rel)
    return slide_wrap(bg, shapes)


# ── Beat → renderer mapping ───────────────────────────────────────────────────

def beat_renderer(beat_type):
    dark_beats = {"hook","moment","pivot","data_reveal","threat","twist"}
    if beat_type in dark_beats:
        return render_hook
    if beat_type == "data_reveal":
        return render_data_reveal
    if beat_type in {"signpost","baseline","map","institution"}:
        return render_section_header
    if beat_type in {"tension","contrast","error_clinic","evaluation_frame"}:
        return render_two_column
    if beat_type == "demonstration":
        return render_diagram
    if beat_type in {"practice","challenge"}:
        return render_question
    if beat_type == "consolidation":
        return render_summary
    return render_content   # mechanism, application, case_study, reveal, etc.


def type_renderer(slide_type):
    """Fallback: map curriculum slide type → renderer."""
    return {
        "title":          render_hook,
        "section_header": render_section_header,
        "content":        render_content,
        "two_column":     render_two_column,
        "diagram":        render_diagram,
        "worked_example": render_worked_example,
        "examiner_tip":   render_examiner_tip,
        "question":       render_question,
        "quote_stat":     render_data_reveal,
        "summary":        render_summary,
    }.get(slide_type, render_content)


# ── PPTX infrastructure ───────────────────────────────────────────────────────

PRES_PROPS_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<p:presentationPr xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'
)
VIEW_PROPS_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<p:viewPr xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
    '<p:normalViewPr><p:restoredLeft sz="15620"/><p:restoredTop sz="94660"/></p:normalViewPr>'
    '</p:viewPr>'
)
TABLE_STYLES_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<a:tblStyleLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
    'def="{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}"/>'
)
SLIDE_MASTER_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
    '  xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
    '  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
    '<p:cSld><p:spTree>'
    '<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
    '<p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/>'
    '<a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>'
    '</p:spTree></p:cSld>'
    '<p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" '
    '  accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" '
    '  accent6="accent6" hlink="hlink" folHlink="folHlink"/>'
    '<p:sldLayoutIdLst><p:sldLayoutId id="2147483649" r:id="rId1"/></p:sldLayoutIdLst>'
    '<p:txStyles><p:titleStyle/><p:bodyStyle/><p:otherStyle/></p:txStyles>'
    '</p:sldMaster>'
)
SLIDE_MASTER_RELS = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
    '<Relationship Id="rId1" '
    'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" '
    'Target="../slideLayouts/slideLayout1.xml"/>'
    '<Relationship Id="rId2" '
    'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" '
    'Target="../theme/theme1.xml"/>'
    '</Relationships>'
)
SLIDE_LAYOUT_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
    '  xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
    '  type="blank" preserve="1">'
    '<p:cSld><p:spTree>'
    '<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
    '<p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/>'
    '<a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>'
    '</p:spTree></p:cSld>'
    '<p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>'
    '</p:sldLayout>'
)
SLIDE_LAYOUT_RELS = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
    '<Relationship Id="rId1" '
    'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" '
    'Target="../slideMasters/slideMaster1.xml"/>'
    '</Relationships>'
)
NOTES_MASTER_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<p:notesMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
    '  xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
    '  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
    '<p:cSld><p:spTree>'
    '<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
    '<p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/>'
    '<a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>'
    '</p:spTree></p:cSld>'
    '<p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" '
    '  accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" '
    '  accent6="accent6" hlink="hlink" folHlink="folHlink"/>'
    '<p:txStyles><p:titleStyle/><p:bodyStyle/><p:otherStyle/></p:txStyles>'
    '</p:notesMaster>'
)
NOTES_MASTER_RELS = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
    '<Relationship Id="rId1" '
    'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" '
    'Target="../theme/theme1.xml"/>'
    '</Relationships>'
)


def theme_xml():
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="Harris Clapham">'
        '<a:themeElements>'
        '<a:clrScheme name="Harris Clapham">'
        f'<a:dk1><a:srgbClr val="{NAVY}"/></a:dk1>'
        f'<a:lt1><a:srgbClr val="{WHITE}"/></a:lt1>'
        f'<a:dk2><a:srgbClr val="{SLATE}"/></a:dk2>'
        f'<a:lt2><a:srgbClr val="{L_GREY}"/></a:lt2>'
        f'<a:accent1><a:srgbClr val="{TEAL}"/></a:accent1>'
        f'<a:accent2><a:srgbClr val="{AMBER}"/></a:accent2>'
        f'<a:accent3><a:srgbClr val="{GREEN}"/></a:accent3>'
        f'<a:accent4><a:srgbClr val="{RED}"/></a:accent4>'
        f'<a:accent5><a:srgbClr val="{TEAL}"/></a:accent5>'
        f'<a:accent6><a:srgbClr val="{SLATE}"/></a:accent6>'
        f'<a:hlink><a:srgbClr val="{TEAL}"/></a:hlink>'
        f'<a:folHlink><a:srgbClr val="{SLATE}"/></a:folHlink>'
        '</a:clrScheme>'
        '<a:fontScheme name="Harris Clapham">'
        '<a:majorFont><a:latin typeface="Calibri Light"/></a:majorFont>'
        '<a:minorFont><a:latin typeface="Calibri"/></a:minorFont>'
        '</a:fontScheme>'
        '<a:fmtScheme name="Office Theme">'
        '<a:fillStyleLst>'
        '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
        '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
        '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
        '</a:fillStyleLst>'
        '<a:lnStyleLst>'
        '<a:ln w="6350"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>'
        '<a:ln w="12700"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>'
        '<a:ln w="19050"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>'
        '</a:lnStyleLst>'
        '<a:effectStyleLst>'
        '<a:effectStyle><a:effectLst/></a:effectStyle>'
        '<a:effectStyle><a:effectLst/></a:effectStyle>'
        '<a:effectStyle><a:effectLst/></a:effectStyle>'
        '</a:effectStyleLst>'
        '<a:bgFillStyleLst>'
        '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
        '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
        '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
        '</a:bgFillStyleLst>'
        '</a:fmtScheme>'
        '</a:themeElements>'
        '</a:theme>'
    )


def notes_slide_xml(notes_text):
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        '  xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        '  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        '<p:cSld><p:spTree>'
        '<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
        '<p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="6858000" cy="9144000"/>'
        '<a:chOff x="0" y="0"/><a:chExt cx="6858000" cy="9144000"/></a:xfrm></p:grpSpPr>'
        '<p:sp><p:nvSpPr><p:cNvPr id="2" name="Slide Image"/>'
        '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
        '<p:nvPr><p:ph type="sldImg"/></p:nvPr></p:nvSpPr><p:spPr/></p:sp>'
        '<p:sp><p:nvSpPr><p:cNvPr id="3" name="Notes"/>'
        '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
        '<p:nvPr><p:ph type="body" sz="quarter" idx="1"/></p:nvPr></p:nvSpPr>'
        '<p:spPr/>'
        '<p:txBody><a:bodyPr/><a:lstStyle/>'
        f'<a:p><a:r><a:rPr lang="en-GB" sz="1000" dirty="0"/>'
        f'<a:t>{xe(str(notes_text))}</a:t></a:r></a:p>'
        '</p:txBody></p:sp>'
        '</p:spTree></p:cSld>'
        '</p:notes>'
    )


def content_types_xml(n_slides, has_logo):
    ov = (
        '<Override PartName="/ppt/theme/theme1.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>\n'
        '<Override PartName="/ppt/slideMasters/slideMaster1.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>\n'
        '<Override PartName="/ppt/slideLayouts/slideLayout1.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>\n'
        '<Override PartName="/ppt/presProps.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.presentationml.presProps+xml"/>\n'
        '<Override PartName="/ppt/viewProps.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.presentationml.viewProps+xml"/>\n'
        '<Override PartName="/ppt/tableStyles.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.presentationml.tableStyles+xml"/>\n'
    )
    for i in range(1, n_slides + 1):
        ov += (f'<Override PartName="/ppt/slides/slide{i}.xml" '
               'ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>\n')
    if has_logo:
        ov += '<Default Extension="png" ContentType="image/png"/>\n'
    ov += '<Default Extension="jpeg" ContentType="image/jpeg"/>\n'
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/ppt/presentation.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>'
        f'{ov}</Types>'
    )


def presentation_xml(n_slides):
    ids = "".join(f'<p:sldId id="{256+i}" r:id="rId{i}"/>\n' for i in range(1, n_slides+1))
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:presentation '
        '  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        '  xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        '  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
        '  saveSubsetFonts="1">'
        '<p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rIdSM"/></p:sldMasterIdLst>'
        '<p:sldSz cx="12192000" cy="6858000" type="screen16x9"/>'
        '<p:notesSz cx="6858000" cy="9144000"/>'
        f'<p:sldIdLst>{ids}</p:sldIdLst>'
        '</p:presentation>'
    )


def presentation_rels_xml(n_slides):
    rels = (
        '<Relationship Id="rIdSM" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" '
        'Target="slideMasters/slideMaster1.xml"/>\n'
        '<Relationship Id="rIdPP" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/presProps" '
        'Target="presProps.xml"/>\n'
        '<Relationship Id="rIdVP" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/viewProps" '
        'Target="viewProps.xml"/>\n'
        '<Relationship Id="rIdTS" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/tableStyles" '
        'Target="tableStyles.xml"/>\n'
    )
    for i in range(1, n_slides + 1):
        rels += (f'<Relationship Id="rId{i}" '
                 'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" '
                 f'Target="slides/slide{i}.xml"/>\n')
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        f'{rels}</Relationships>'
    )


def slide_rels_xml(idx, img_map, has_notes):
    rels = (
        '<Relationship Id="rId1" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" '
        'Target="../slideLayouts/slideLayout1.xml"/>\n'
    )
    for rid, target in sorted(img_map.items()):
        rels += (f'<Relationship Id="{rid}" '
                 'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" '
                 f'Target="{target}"/>\n')
    if has_notes:
        rels += (f'<Relationship Id="rIdN" '
                 'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide" '
                 f'Target="../notesSlides/notesSlide{idx}.xml"/>\n')
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        f'{rels}</Relationships>'
    )


def notes_rels_xml(idx):
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        f'<Relationship Id="rId1" '
        f'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" '
        f'Target="../slides/slide{idx}.xml"/>'
        f'<Relationship Id="rId2" '
        f'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster" '
        f'Target="../notesMasters/notesMaster1.xml"/>'
        '</Relationships>'
    )


# ── Namespace constants ───────────────────────────────────────────────────────

_P  = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A  = "http://schemas.openxmlformats.org/drawingml/2006/main"
_R  = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _add_image_rel(slide, abs_path):
    """Add an image to a slide via python-pptx and return its relationship ID."""
    with open(abs_path, "rb") as f:
        bio = BytesIO(f.read())
    ph = slide.shapes.add_picture(bio, 0, 0)
    blip = ph._element.find(f".//{{{_A}}}blip")
    rId = blip.get(f"{{{_R}}}embed")
    ph._element.getparent().remove(ph._element)
    return rId


def _inject_slide(slide, slide_xml):
    """Replace a python-pptx slide's cSld content with our rendered XML."""
    src = etree.fromstring(slide_xml.encode("utf-8"))
    our_cSld   = src.find(f"{{{_P}}}cSld")
    our_bg     = our_cSld.find(f"{{{_P}}}bg")
    our_spTree = our_cSld.find(f"{{{_P}}}spTree")

    slide_cSld   = slide._element.find(f"{{{_P}}}cSld")
    slide_spTree = slide_cSld.find(f"{{{_P}}}spTree")

    # Insert background before spTree
    if our_bg is not None:
        slide_cSld.insert(list(slide_cSld).index(slide_spTree), our_bg)

    # Clear default placeholder shapes (keep first 2: nvGrpSpPr + grpSpPr)
    for child in list(slide_spTree)[2:]:
        slide_spTree.remove(child)

    # Append our shapes
    for child in list(our_spTree)[2:]:
        slide_spTree.append(child)


# ── Main build ────────────────────────────────────────────────────────────────

def build(unit_path, story_path, output_path):
    unit_path  = Path(unit_path)
    story_path = Path(story_path)
    base_dir   = unit_path.parent

    with open(unit_path)  as f: unit  = json.load(f)
    with open(story_path) as f: story = json.load(f)

    # Index story briefs by slide_ref
    briefs = {b["slide_ref"]: b for b in story["slides"]}

    # Flatten curriculum slides
    all_slides = []
    for section in unit["sections"]:
        for sl in section.get("slides", []):
            all_slides.append(sl)

    # Logo
    logo_rel_path = story.get("meta",{}).get("style",{}).get("logo","")
    logo_abs = base_dir / logo_rel_path if logo_rel_path else None
    has_logo = bool(logo_abs and logo_abs.exists())

    # ── Phase 1: build infrastructure via python-pptx, collect rIds ──────────
    import zipfile as _zf
    prs = Presentation()
    prs.slide_width  = Emu(W)
    prs.slide_height = Emu(H)
    blank_layout = prs.slide_layouts[6]

    slide_render_data = []   # (brief, sl, logo_rel_id, img_rels)

    for sl in all_slides:
        slide = prs.slides.add_slide(blank_layout)
        logo_rel_id = _add_image_rel(slide, logo_abs) if has_logo else None
        img_rels = {}
        for k in ("image", "diagram_image"):
            if sl.get(k):
                abs_p = base_dir / sl[k]
                if abs_p.exists():
                    img_rels[sl[k]] = _add_image_rel(slide, abs_p)
        slide_render_data.append((sl, logo_rel_id, img_rels))

    # ── Phase 2: save to buffer, then patch slide XML files in the ZIP ────────
    buf = BytesIO()
    prs.save(buf)

    import re as _re
    buf.seek(0)
    out = Path(output_path)
    with _zf.ZipFile(buf) as src_zip, _zf.ZipFile(out, "w", _zf.ZIP_DEFLATED) as dst_zip:
        for name in src_zip.namelist():
            m = _re.match(r"ppt/slides/slide(\d+)\.xml$", name)
            if m:
                idx = int(m.group(1))
                sl, logo_rel_id, img_rels = slide_render_data[idx - 1]
                brief   = briefs.get(sl.get("slide_id", ""), {})
                beat    = brief.get("beat_type", "")
                sl_type = sl.get("type", "content")
                renderer = beat_renderer(beat) if beat else type_renderer(sl_type)
                if sl_type == "examiner_tip":              renderer = render_examiner_tip
                elif sl_type == "question":                renderer = render_question
                elif sl_type == "worked_example":          renderer = render_worked_example
                elif sl_type == "quote_stat":              renderer = render_data_reveal
                elif sl_type == "two_column" and not beat: renderer = render_two_column
                elif sl_type == "diagram":                 renderer = render_diagram
                reset_ids()
                dst_zip.writestr(name, renderer(brief, sl, logo_rel_id, img_rels).encode("utf-8"))
            else:
                dst_zip.writestr(name, src_zip.read(name))

    print(f"✅ {out}  ({len(all_slides)} slides)")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python3 render.py <curriculum_unit.json> <story.json> [output.pptx]")
        sys.exit(1)
    unit_p  = sys.argv[1]
    story_p = sys.argv[2]
    out_p   = sys.argv[3] if len(sys.argv) > 3 else str(Path(sys.argv[1]).parent / "Protectionism.pptx")
    build(unit_p, story_p, out_p)
